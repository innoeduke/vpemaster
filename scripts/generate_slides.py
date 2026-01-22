import os
import pymysql
import re
from datetime import datetime
from dotenv import load_dotenv
from pptx import Presentation

load_dotenv()

def info_fmt(name, creds):
    if not name:
        return ""
    if not creds:
        return name
    return f"{name}, {creds}"

def fetch_meeting_data(meeting_num):
    db_url = os.getenv("DATABASE_URL")
    if not db_url:
        print("Error: DATABASE_URL not found in .env")
        return None
    
    # Extract DB info from URL (mysql+pymysql://user:pass@host/db)
    match = re.match(r"mysql\+pymysql://([^:]+):([^@]+)@([^/]+)/(.+)", db_url)
    if not match:
        print("Error: Could not parse DATABASE_URL")
        return None
        
    user, password, host, db_name = match.groups()
    
    conn = pymysql.connect(
        host=host,
        user=user,
        password=password,
        database=db_name,
        cursorclass=pymysql.cursors.DictCursor
    )
    
    data = {}
    try:
        with conn.cursor() as cursor:
            # 1. Get Meeting and Club info
            sql = """
            SELECT m.Meeting_Date, m.Meeting_Number, c.club_name 
            FROM Meetings m
            JOIN clubs c ON m.club_id = c.id
            WHERE m.Meeting_Number = %s
            """
            cursor.execute(sql, (meeting_num,))
            meeting = cursor.fetchone()
            if not meeting:
                print(f"Error: Meeting {meeting_num} not found.")
                return None
            
            data['meeting_number'] = str(meeting['Meeting_Number'])
            if meeting['Meeting_Date']:
                data['meeting_date'] = meeting['Meeting_Date'].strftime("%d-%b-%Y")
            data['club_name'] = meeting['club_name']
            
            # 2. Get Session Logs
            sql = """
            SELECT 
                sl.id, 
                sl.Meeting_Seq, 
                st.Title as Session_Type, 
                mr.name as Role_Name, 
                con.Name as Owner_Name,
                sl.credentials,
                sl.Session_Title,
                p.Project_Name,
                sl.project_code
            FROM Session_Logs sl
            JOIN Session_Types st ON sl.Type_ID = st.id
            LEFT JOIN meeting_roles mr ON st.role_id = mr.id
            LEFT JOIN Contacts con ON sl.Owner_ID = con.id
            LEFT JOIN Projects p ON sl.Project_ID = p.id
            WHERE sl.Meeting_Number = %s
            ORDER BY sl.Meeting_Seq
            """
            cursor.execute(sql, (meeting_num,))
            logs = cursor.fetchall()
            
            # Helper to find first non-None owner for a role or session title part
            def find_role_info(role_name_pattern=None, session_title_pattern=None):
                for l in logs:
                    match_role = role_name_pattern and l['Role_Name'] and re.search(role_name_pattern, l['Role_Name'], re.I)
                    match_title = session_title_pattern and l['Session_Title'] and re.search(session_title_pattern, l['Session_Title'], re.I)
                    if (match_role or match_title) and l['Owner_Name']:
                        return info_fmt(l['Owner_Name'], l['credentials'])
                return ""

            # Standard Roles
            data['saa_info'] = find_role_info(session_title_pattern="SAA Intro")
            data['president_info'] = find_role_info(session_title_pattern="President's Address")
            data['welcome-officer_info'] = find_role_info(role_name_pattern="Welcome Officer")
            data['tme_info'] = find_role_info(role_name_pattern="Toastmaster")
            data['timer_info'] = find_role_info(role_name_pattern="Timer")
            data['ah-counter_info'] = find_role_info(role_name_pattern="Ah-Counter")
            data['grammarian_info'] = find_role_info(role_name_pattern="Grammarian")
            data['topicsmaster_info'] = find_role_info(role_name_pattern="Topicsmaster")
            data['ge_info'] = find_role_info(role_name_pattern="General Evaluator")
            
            # Speeches
            speakers = [l for l in logs if l['Role_Name'] == "Prepared Speaker" and l['Owner_Name']]
            for i in range(1, 4):
                if len(speakers) >= i:
                    s = speakers[i-1]
                    data[f'ps{i}-speaker'] = s['Owner_Name']
                    data[f'ps{i}_title'] = s['Project_Name'] or s['Session_Title'] or ""
                    data[f'ps{i}-project_info'] = f"{s['project_code']} {s['Project_Name']}" if s['project_code'] and s['Project_Name'] else (s['Project_Name'] or "")
                    data[f'ps{i}-speaker_info'] = info_fmt(s['Owner_Name'], s['credentials'])
                else:
                    data[f'ps{i}-speaker'] = ""
                    data[f'ps{i}_title'] = ""
                    data[f'ps{i}-project_info'] = ""
                    data[f'ps{i}-speaker_info'] = ""
            
            # Evaluators
            evaluators = [l for l in logs if l['Role_Name'] == "Individual Evaluator" and l['Owner_Name']]
            for i in range(1, 4):
                if len(evaluators) >= i:
                    e = evaluators[i-1]
                    data[f'ie{i}_info'] = info_fmt(e['Owner_Name'], e['credentials'])
                else:
                    data[f'ie{i}_info'] = ""

            # Keynote
            data['keynote_title'] = find_role_info(role_name_pattern="Keynote Speaker") # Usually just title? User extracted keynote_title.
            # Let's adjust keynote
            k = next((l for l in logs if l['Role_Name'] == "Keynote Speaker"), None)
            if k:
                data['keynote_title'] = k['Session_Title'] or k['Project_Name'] or ""
            else:
                data['keynote_title'] = ""

            # Placeholder fallback for missing ones
            all_vars = [
                'saa_info', 'club_name', 'welcome-officer_info', 'president_info', 
                'vpm_info', 'tme_info', 'timer_info', 'ah-counter_info', 
                'grammarian_info', 'topicsmaster_info', 'ge_info', 
                'ps1-speaker', 'ps1_title', 'ps1-project_info', 'ps1-speaker_info',
                'ps2-speaker', 'ps2_title', 'ps2-project_info', 'ps2-speaker_info',
                'ps3-speaker', 'ps3_title', 'ps3-project_info', 'ps3-speaker_info',
                'ie1_info', 'ie2_info', 'ie3_info', 'keynote_title'
            ]
            for v in all_vars:
                if v not in data or not data[v]:
                    data[v] = "" # or "TBD"
                    
    finally:
        conn.close()
    
    return data

def replace_placeholders(pptx_path, replacements, output_path):
    if not os.path.exists(pptx_path):
        print(f"Error: File {pptx_path} not found.")
        return
    # Load the presentation
    prs = Presentation(pptx_path)
    
    # Process replacements to include double braces if they don't have them
    final_replacements = {}
    for k, v in replacements.items():
        key = k if k.startswith("{{") else f"{{{{{k}}}}}"
        final_replacements[key] = v

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            for paragraph in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in paragraph.runs)
                
                updated_text = full_text
                has_change = False
                for key, value in final_replacements.items():
                    if key in updated_text:
                        updated_text = updated_text.replace(key, str(value))
                        has_change = True
                
                if has_change:
                    if paragraph.runs:
                        paragraph.runs[0].text = updated_text
                        for i in range(1, len(paragraph.runs)):
                            paragraph.runs[i].text = ""
                    else:
                        paragraph.text = updated_text
    # Save the updated presentation
    prs.save(output_path)
    print(f"Successfully saved updated presentation to: {output_path}")

if __name__ == "__main__":
    meeting_number = 959
    
    # Define paths
    home = os.path.expanduser("~")
    # Template location
    # Based on earlier list_dir, it's in /Users/wmu/workspace/toastmasters/vpemaster/instance/
    base_dir = "/Users/wmu/workspace/toastmasters/vpemaster"
    input_file = os.path.join(base_dir, "instance", "SHLTMC_Meeting_<nnn>.pptx")
    output_file = os.path.join(home, "Downloads", f"SHLTMC_Meeting_{meeting_number}.pptx")
    
    # Fetch data
    data_to_replace = fetch_meeting_data(meeting_number)
    
    if data_to_replace:
        replace_placeholders(input_file, data_to_replace, output_file)
    else:
        print("Failed to fetch meeting data.")
