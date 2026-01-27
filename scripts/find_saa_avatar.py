from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import copy
from app import create_app, db
from app.models import Club
from PIL import Image

def crop_image_to_aspect_ratio(image_path, target_width, target_height):
    """
    Crops an image to match the target aspect ratio (center crop).
    Returns path to cropped image.
    """
    try:
        im = Image.open(image_path)
        img_w, img_h = im.size
        
        target_ratio = target_width / target_height
        img_ratio = img_w / img_h
        
        if img_ratio > target_ratio:
            # Image is wider: crop width
            new_width = int(img_h * target_ratio)
            left = (img_w - new_width) // 2
            box = (left, 0, left + new_width, img_h)
        else:
            # Image is taller: crop height
            new_height = int(img_w / target_ratio)
            top = (img_h - new_height) // 2
            box = (0, top, img_w, top + new_height)
            
        cropped_im = im.crop(box)
        
        # Convert to RGB for JPEG compatibility
        if cropped_im.mode in ('RGBA', 'LA'):
            background = Image.new(cropped_im.mode[:-1], cropped_im.size, (255, 255, 255))
            background.paste(cropped_im, mask=cropped_im.split()[-1])
            cropped_im = background.convert('RGB')
        elif cropped_im.mode != 'RGB':
            cropped_im = cropped_im.convert('RGB')
        
        base, ext = os.path.splitext(image_path)
        cropped_path = f"{base}_cropped.jpg"
        cropped_im.save(cropped_path, 'JPEG')
        return cropped_path
    except Exception as e:
        print(f"Error cropping image: {e}")
        return None

def apply_shape_geometry(picture, source_shape):
    """
    Apply the geometry (prstGeom) of the source_shape to the picture.
    """
    try:
        # Define namespaces
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        
        # Access XML elements
        pic_element = picture._element
        shape_element = source_shape._element
        
        # Helper to find child with namespace
        def find_child(element, ns, tag_name):
            return element.find(f"{{{ns}}}{tag_name}")
            
        # Source spPr (PresentationML)
        source_spPr = find_child(shape_element, ns_p, 'spPr')
        if source_spPr is None:
            # Fallback: maybe it's not a generic shape, but let's try
            source_spPr = find_child(shape_element, ns_a, 'spPr')
        
        if source_spPr is None:
            print(f"Source shape has no spPr. Tag: {shape_element.tag}")
            return
            
        # Source prstGeom (DrawingML)
        source_prstGeom = find_child(source_spPr, ns_a, 'prstGeom')
        if source_prstGeom is None:
            print("Source shape has no preset geometry.")
            return

        # Target spPr (PresentationML)
        target_spPr = find_child(pic_element, ns_p, 'spPr')
        if target_spPr is None:
            print("Target picture has no spPr.")
            return
            
        # Target prstGeom (DrawingML)
        target_prstGeom = find_child(target_spPr, ns_a, 'prstGeom')
        
        # Replace target geometry with copy of source geometry
        if target_prstGeom is not None:
             target_spPr.remove(target_prstGeom)
             
        target_spPr.append(copy.deepcopy(source_prstGeom))
        # print(f"Applied geometry {source_prstGeom.attrib} to picture.")

    except Exception as e:
        print(f"Error applying geometry: {e}")

def replace_saa_avatar(pptx_path, shape_name, output_path, club_id=1):
    # 1. Lookup SAA Avatar URL
    app = create_app()
    with app.app_context():
        club = db.session.get(Club, club_id)
        if not club or not club.current_excomm or not club.current_excomm.saa:
            print("SAA or Avatar URL not found.")
            return
        
        saa = club.current_excomm.saa
        if not saa.Avatar_URL:
            print(f"No avatar URL for SAA {saa.Name}")
            return
            
        print(f"Found SAA: {saa.Name}, Avatar: {saa.Avatar_URL}")
        image_path = os.path.join(app.static_folder, saa.Avatar_URL)
        if not os.path.exists(image_path):
            print(f"Image file not found at: {image_path}")
            return

    # 2. Replace Shape in PPTX
    if not os.path.exists(pptx_path):
        print(f"PPTX file not found: {pptx_path}")
        return

    prs = Presentation(pptx_path)
    replaced = False

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.name == shape_name:
                print(f"Found target shape on Slide {i+1}")
                
                # Capture properties
                left, top = shape.left, shape.top
                width, height = shape.width, shape.height
                
                # Crop image
                cropped_path = crop_image_to_aspect_ratio(image_path, width, height)
                if not cropped_path:
                    break
                    
                try:
                    # Add new picture
                    pic = slide.shapes.add_picture(cropped_path, left, top, width, height)
                    
                    # Apply geometry from old shape to new picture
                    apply_shape_geometry(pic, shape)
                    
                    # Remove original shape
                    sp = shape.element
                    sp.getparent().remove(sp)
                    
                    print("Replaced shape with cropped and shaped image.")
                    replaced = True
                    
                    # Cleanup
                    try: os.remove(cropped_path)
                    except: pass
                    
                except Exception as e:
                    print(f"Error replacing shape: {e}")
                
                break 
        if replaced: 
            break

    if replaced:
        prs.save(output_path)
        print(f"Saved modified presentation to: {output_path}")
    else:
        print("Target shape not found or replacement failed.")

if __name__ == "__main__":
    pptx_file = "instance/test_slides.pptx"
    output_file = "instance/test_slides_with_avatar.pptx"
    target_shape = "saa_avatar"
    replace_saa_avatar(pptx_file, target_shape, output_file)
