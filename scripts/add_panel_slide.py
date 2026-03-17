from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
import copy

template_path = '/Users/wmu/workspace/toastmasters/vpemaster/app/static/club_resources/2/slides_template.pptx'
prs = Presentation(template_path)

# Clone Slide 30 (IE slide) to get the same background/style
source_slide = prs.slides[29] # 0-indexed, slide 30
dest_slide = prs.slides.add_slide(source_slide.slide_layout)

# Remove all shapes from the newly created slide to start fresh (but keep the background)
for shape in dest_slide.shapes:
    # Most themes handle background in the layout, so we just clear specific shapes
    pass

# Helper to add a person (avatar + info)
def add_person(slide, avatar_name, info_placeholder, left, top, avatar_size, fontSize=24):
    # Add avatar placeholder (as a shape named correctly)
    avatar_shape = slide.shapes.add_shape(
        1, # 1 is Rect
        Inches(left), Inches(top), Inches(avatar_size), Inches(avatar_size)
    )
    avatar_shape.name = avatar_name
    
    # Add text placeholder
    text_box = slide.shapes.add_textbox(
        Inches(left), Inches(top + avatar_size + 0.1), Inches(avatar_size * 1.5), Inches(0.5)
    )
    text_frame = text_box.text_frame
    text_frame.text = info_placeholder
    p = text_frame.paragraphs[0]
    p.font.size = Pt(fontSize)
    p.alignment = PP_ALIGN.CENTER

# Clear destination slide (except for items we want to keep, like background if any)
# Actually, add_slide creates a new slide with layout. Let's just use it.
for shape in list(dest_slide.shapes):
    if shape.name != "Title 4": # Keep title if it's there? No, let's just clear.
        sp = shape._element
        sp.getparent().remove(sp)

# Add Title
title_box = dest_slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "PANEL DISCUSSION"
p = title_frame.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Layout dimensions
slide_width = 13.33
slide_height = 7.5

# Moderator (Top Centerish)
mod_size = 2.0
add_person(dest_slide, "moderator_avatar", "{{moderator_info}}", (slide_width - mod_size)/2, 1.8, mod_size, 18)

# Panelists (Row below)
pan_size = 1.6
spacing = 0.5
total_width = (pan_size * 4) + (spacing * 3)
start_x = (slide_width - total_width) / 2
start_y = 4.5

for i in range(1, 5):
    add_person(dest_slide, f"panelist{i}_avatar", f"{{{{panelist{i}_info}}}}", start_x + (i-1)*(pan_size + spacing), start_y, pan_size, 14)

# Save the updated template
prs.save(template_path)
print(f"Successfully updated template with Panel Discussion slide at {template_path}")
