import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_slide(prs, title, content_points):
    """Helper to create a bullet point slide."""
    slide_layout = prs.slide_layouts[1]  # Bullet layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Content
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Clear existing default text

    for point in content_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(24)

    return slide

def create_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.shapes.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle

def main():
    prs = Presentation()

    # 1. Title Slide
    create_title_slide(prs, "VPEMaster User Guide", "Getting Started with Meeting Manager")

    # 2. Introduction
    create_slide(prs, "Introduction", [
        "Welcome to VPEMaster/Meeting Manager!",
        "This tool helps Toastmasters clubs manage:",
        " - Meeting Agendas",
        " - Member Rosters & Contacts",
        " - User Accounts & Roles",
        " - Speech Progress Tracking"
    ])

    # 3. Logging In
    create_slide(prs, "Logging In", [
        "Access the application via your browser.",
        "Default Password for new users: 'leadership'",
        "Note: You will be prompted to change your password after logging in (recommended)."
    ])

    # 4. Managing Contacts
    create_slide(prs, "Contact Book", [
        "Navigate to 'Contacts' to view all people.",
        "Features:",
        " - Search by Name",
        " - Add New Contacts (Guests/Members)",
        " - View Club Affiliation and Education Status"
    ])

    # 5. Quick Add User (New Feature)
    create_slide(prs, "Quick Add User Feature", [
        "Easily create login accounts for existing contacts.",
        "1. Go to the Contact Book.",
        "2. Look for the 'Person Plus' icon in the Actions column.",
        "3. Click to auto-create a user.",
        "   - Auto-generated username (e.g., 'chrislee').",
        "   - Default password: 'leadership'.",
        "   - Icon is grayed out if user already exists."
    ])

    # 6. Agenda Management
    create_slide(prs, "Agenda Management", [
        "View upcoming meetings from the Dashboard.",
        "Assign roles (Speakers, Evaluators, etc.).",
        "Export Agenda to Excel:",
        " - Standard Agenda Sheet",
        " - Participants Sheet: Grouped list of all attendees for easy tracking."
    ])

    # 7. User Profile & Security
    create_slide(prs, "Profile & Security", [
        "Click your name in the top right to access Profile.",
        "Update your bio and contact info.",
        "Change Password:",
        " - Must be 8+ characters.",
        " - Requires Uppercase, Lowercase, and Number.",
        " - Note: Admin created users bypass this for the initial 'leadership' password."
    ])

    output_file = "VPEMaster_User_Guide.pptx"
    prs.save(output_file)
    print(f"Successfully generated {output_file}")

if __name__ == "__main__":
    main()
