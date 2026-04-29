import io
import os
import re
import copy
import traceback
import math
from flask import current_app
from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

from .export.context import MeetingExportContext
from ..utils import derive_credentials
from ..models.meeting import Meeting
from ..models.contact import Contact
from .. import db

class MeetingSlideService:
    """Service to generate meeting PowerPoint slides."""

    @staticmethod
    def generate_meeting_pptx(meeting_id, logs_data=None):
        """Main entry point for generating meeting slides. Uses v2 by default."""
        return MeetingSlideService.generate_meeting_pptx_v2(meeting_id, logs_data)

    @staticmethod
    def generate_meeting_pptx_v2(meeting_id, logs_data=None):
        """
        Version 2 (Layout-based) slide generation.
        Uses slides_layouts.pptx from club_resources.
        """
        meeting = db.session.get(Meeting, meeting_id)
        if not meeting:
            current_app.logger.error(f"Meeting {meeting_id} not found.")
            return None

        # If logs_data not provided, fetch it (fallback, though preferred to pass it from route)
        if logs_data is None:
            from ..agenda_routes import _get_processed_logs_data
            logs_data, _ = _get_processed_logs_data(meeting_id)

        # Template Path: app/static/club_resources/<club_id>/slides_layouts.pptx
        template_path = os.path.join(current_app.static_folder, 'club_resources', str(meeting.club_id), 'slides_layouts.pptx')
        
        # Fallback to generic template if club-specific one doesn't exist
        if not os.path.exists(template_path):
             current_app.logger.warning(f"Club-specific layout template not found at {template_path}. Falling back to instance/layouts.pptx")
             template_path = os.path.join(current_app.root_path, '..', 'instance', 'layouts.pptx')
        
        if not os.path.exists(template_path):
            current_app.logger.error(f"Template not found at {template_path}")
            return None

        try:
            prs = Presentation(template_path)
            
            # Available layouts mapping for easy lookup
            layouts = {layout.name: layout for layout in prs.slide_layouts}

            # --- 1. Add Title Slide ---
            title_layout = layouts.get('Title Slide')
            if title_layout:
                title_slide = prs.slides.add_slide(title_layout)
                club_name = (meeting.club.club_name if meeting.club else "Toastmasters Club")
                meeting_date_str = meeting.Meeting_Date.strftime('%d-%b-%Y') if meeting.Meeting_Date else ""
                meeting_info = f"Meeting {meeting.Meeting_Number} / {meeting_date_str}"
                
                for shape in title_slide.placeholders:
                    idx = shape.placeholder_format.idx
                    if idx == 0: shape.text = club_name
                    elif idx == 1: shape.text = meeting_info
            
            # --- 2. Add Action/Agenda Slide (right after Title Slide) ---
            action_layout = layouts.get('section_action') or layouts.get('section_agenda')
            if action_layout:
                prs.slides.add_slide(action_layout)

            avatar_base_path = os.path.join(current_app.static_folder, 'uploads', 'avatars')

            for log in logs_data:
                session_title = (log.get('Session_Title') or "").strip()
                session_type = (log.get('session_type_title') or "").strip()
                is_section = log.get('is_section')
                
                # Normalize title for matching
                title_upper = session_title.upper()
                
                # --- 3. Section detection based on Session Title ---
                section_layout = None
                if "OPENING" in title_upper:
                    section_layout = layouts.get('section_opening')
                elif ("EVALUATION" in title_upper or "EVALUATOR" in title_upper) and is_section:
                    section_layout = layouts.get('section_evaluations')
                elif "PREPARED SPEECHES" in title_upper:
                    section_layout = layouts.get('section_preparedspeeches')
                elif "AWARDS & CLOSING" in title_upper:
                    section_layout = layouts.get('section_voting')
                
                if section_layout:
                    prs.slides.add_slide(section_layout)
                    continue

                # Durations and visibility check
                dur_min = log.get('Duration_Min') or 0
                dur_max = log.get('Duration_Max') or 0
                if dur_max == 0:
                    continue

                # Standard Role Slide Selection
                target_layout = None
                if session_type == 'Networking':
                    target_layout = layouts.get('section_networking')
                elif session_type == 'Keynote Speech':
                    target_layout = layouts.get('Keynote Speaker Slide')
                elif session_type == 'Prepared Speech':
                    target_layout = layouts.get('Prepared Speaker Slide')
                else:
                    target_layout = layouts.get('Role Taker Slide')
                
                if not target_layout:
                    target_layout = layouts.get('Role Taker Slide')
                if not target_layout:
                    continue

                slide = prs.slides.add_slide(target_layout)
                
                # Content preparation
                if session_type == 'Evaluation':
                    title_text = f"Individual Evaluator for\n{session_title}"
                else:
                    title_text = session_title or session_type or 'Untitled'

                name = log.get('owner_name') or ''
                creds = log.get('Credentials') or ''
                subtitle_text = f"{name}, {creds}" if creds and name else (name or creds)
                duration_text = f"{dur_min}-{dur_max}'" if dur_min != dur_max else f"{dur_max}'"

                project_info_text = ""
                if session_type == 'Prepared Speech':
                    p_code = log.get('project_code') or ""
                    p_name = log.get('project_name') or ""
                    project_info_text = f"{p_code} - {p_name}" if p_code and p_name else (p_code or p_name)

                # Placeholder assignment
                for shape in slide.placeholders:
                    idx = shape.placeholder_format.idx
                    if idx == 0: shape.text = title_text
                    elif idx == 1: shape.text = subtitle_text
                    elif idx == 10: shape.text = duration_text
                    elif idx == 12: shape.text = project_info_text
                    elif idx == 11:
                        # Avatar insertion with WEBP conversion
                        primary_owner_id = log.get('Owner_ID')
                        if primary_owner_id:
                            contact = Contact.query.get(primary_owner_id)
                            if contact and contact.Avatar_URL:
                                avatar_path = os.path.join(avatar_base_path, contact.Avatar_URL)
                                if os.path.exists(avatar_path):
                                    try:
                                        # Process all images to ensure transparency is filled with white
                                        img = Image.open(avatar_path)
                                        if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
                                            background = Image.new('RGB', img.size, (255, 255, 255))
                                            if img.mode == 'P':
                                                img = img.convert('RGBA')
                                            background.paste(img, mask=img.split()[-1])
                                            img = background
                                        elif img.mode != 'RGB':
                                            img = img.convert('RGB')
                                        
                                        img_io = io.BytesIO()
                                        img.save(img_io, format='PNG')
                                        img_io.seek(0)
                                        shape.insert_picture(img_io)
                                    except Exception as e:
                                        current_app.logger.error(f"Error inserting picture {avatar_path}: {e}")
                
                # Table Topics divider
                if session_type == 'Table Topics':
                    tt_divider_layout = layouts.get('section_tabletopics')
                    if tt_divider_layout:
                        prs.slides.add_slide(tt_divider_layout)
            
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            return output

        except Exception as e:
            current_app.logger.error(f"Error generating PPTX v2: {e}\n{traceback.format_exc()}")
            return None

    # --- Legacy v1 Methods ---

    @staticmethod
    def generate_meeting_pptx_v1(meeting_id):
        """Legacy v1 (Template-based replacement) slide generation."""
        context = MeetingExportContext(meeting_id)
        meeting = context.meeting
        if not meeting:
            return None

        def info_fmt(name, creds, role=None):
            if not name: return ""
            full_name = f"{name}, {creds}" if creds else name
            return f"{full_name}\n{role}" if role else full_name

        def dur_fmt(log):
            if not log: return ""
            dmin, dmax = log.Duration_Min, log.Duration_Max
            if dmin is not None and dmax is not None:
                if dmin == 0: return f"{dmax} '"
                return f"{dmin} ~ {dmax} '"
            val = dmax if dmax is not None else dmin
            return f"{val} '" if val is not None else ""

        replacements = MeetingSlideService._initialize_placeholders(meeting)
        avatar_map = {
            "saa_avatar": None, "welcome-officer_avatar": None, "president_avatar": None,
            "vpm_avatar": None, "vpe_avatar": None, "vppr_avatar": None,
            "treasurer_avatar": None, "secretary_avatar": None, "tme_avatar": None,
            "timer_avatar": None, "ah-counter_avatar": None, "grammarian_avatar": None,
            "topicsmaster_avatar": None, "ge_avatar": None, "photographer_avatar": None,
            "keynote-speaker_avatar": None, "moderator_avatar": None
        }
        for i in range(1, 7):
            avatar_map[f"ps{i}_avatar"] = None
            avatar_map[f"ie{i}_avatar"] = None
        for i in range(1, 5):
            avatar_map[f"panelist{i}_avatar"] = None
        
        MeetingSlideService._populate_standard_roles(context, meeting, replacements, info_fmt, dur_fmt, avatar_map)
        MeetingSlideService._populate_speakers(context, replacements, info_fmt, dur_fmt, avatar_map)
        MeetingSlideService._populate_evaluators(context, replacements, info_fmt, dur_fmt, avatar_map)
        MeetingSlideService._populate_featured_session(context, meeting, replacements, info_fmt, dur_fmt, avatar_map)
        MeetingSlideService._populate_table_topics(context, replacements, dur_fmt)

        template_path = os.path.join(current_app.static_folder, 'club_resources', str(meeting.club_id), 'slides_template.pptx')
        if not os.path.exists(template_path):
            current_app.logger.error(f"PPTX Template not found at: {template_path}")
            return None

        try:
            prs = Presentation(template_path)
            for slide in prs.slides:
                host_type = MeetingSlideService._identify_host_type(slide)
                MeetingSlideService._populate_slide(slide, meeting, context, replacements, avatar_map, host_type)
            
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            return output
        except Exception as e:
            current_app.logger.error(f"Error generating PPTX v1: {e}\n{traceback.format_exc()}")
            return None

    @staticmethod
    def add_event_template_slide(club_id, slide_title, host, extra_contacts, duration_info, event_type=None, slide_number=None):
        """Add a template slide for a specific event (Legacy)."""
        template_path = os.path.join(current_app.static_folder, 'club_resources', str(club_id), 'slides_template.pptx')
        if not os.path.exists(template_path): return False
        try:
            prs = Presentation(template_path)
            layout = next((l for l in prs.slide_layouts if l.name == "Title Only"), prs.slide_layouts[0])
            slide = prs.slides.add_slide(layout)
            if slide_number is not None:
                idx = max(0, min(slide_number - 1, len(prs.slides) - 1))
                xml_slides = prs.slides._sldIdLst
                slide_id = xml_slides[-1]
                xml_slides.remove(slide_id)
                xml_slides.insert(idx, slide_id)

            title_shape = slide.shapes.title
            if title_shape:
                role = host.get('role', 'host').lower().replace(' ', '_')
                title_shape.text = f"{{{{{role}_session_title}}}}".upper() if not slide_title else slide_title.upper()
                if title_shape.has_text_frame:
                    p = title_shape.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    if p.runs:
                        p.runs[0].font.name = "Gotham"
                        p.runs[0].font.size = Pt(40.5)
                        p.runs[0].font.color.rgb = RGBColor.from_string('F2DF74')

            font_name, tm_red = "Gotham", RGBColor(228, 31, 53)
            dur_box = slide.shapes.add_textbox(Cm(1.0), Cm(1.0), Cm(4.23), Cm(1.16))
            tf = dur_box.text_frame
            tf.text = duration_info if duration_info and "{{" not in duration_info else "{{slide_duration}}"
            tf.vertical_anchor = 3
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.name, p.runs[0].font.size = font_name, Pt(12)
            p.runs[0].font.color.rgb = RGBColor(255, 255, 255)

            host_pic = slide.shapes.add_shape(9, Cm(2.5), Cm(7.0), Cm(5.0), Cm(5.0))
            host_pic.name = "host_avatar"
            host_pic.line.color.rgb, host_pic.line.width = RGBColor(119, 36, 50), Pt(3)

            prs.save(template_path)
            return True
        except Exception as e:
            current_app.logger.error(f"Error adding event template slide: {e}")
            return False

    @staticmethod
    def _initialize_placeholders(meeting):
        """Internal v1 helper."""
        replacements = {
            "{{meeting_number}}": str(meeting.Meeting_Number),
            "{{meeting_date}}": meeting.Meeting_Date.strftime('%d-%b-%Y') if meeting.Meeting_Date else "",
            "{{club_name}}": (meeting.club.club_name if meeting.club else ""),
        }
        for p in ["saa", "welcome-officer", "president", "vpm", "vpe", "vppr", "treasurer", "secretary", "tme", "timer", "ah-counter", "grammarian", "topicsmaster", "ge", "photographer"]:
            replacements.update({f"{{{{{p}_info}}}}": "", f"{{{{{p}_duration}}}}": "", f"{{{{{p}_session_title}}}}": ""})
        for i in range(1, 7):
            replacements.update({f"{{{{ps{i}}}}}": "", f"{{{{ps{i}_title}}}}": "", f"{{{{ps{i}-project_info}}}}": "", f"{{{{ps{i}_info}}}}": "", f"{{{{ps{i}_duration}}}}": "", f"{{{{ps{i}_session_title}}}}": ""})
            replacements.update({f"{{{{ie{i}_info}}}}": "", f"{{{{ie{i}_duration}}}}": "", f"{{{{ie{i}_session_title}}}}": ""})
        replacements.update({"{{keynote_title}}": "", "{{keynote_duration}}": "", "{{keynote-speaker_info}}": "", "{{keynote_session_title}}": "", "{{table-topics_duration}}": "", "{{table-topics_session_title}}": ""})
        return replacements

    @staticmethod
    def _populate_standard_roles(context, meeting, replacements, info_fmt, dur_fmt, avatar_map=None):
        """Internal v1 helper."""
        def populate_role(role_pattern=None, title_pattern=None, info_key=None, dur_key=None):
            last_match = None
            for log, st in context.logs:
                if not log.owner: continue
                role_name = (st.role.name if st.role else st.Title)
                if (role_pattern and role_name and re.search(role_pattern, role_name, re.I)) or (title_pattern and log.Session_Title and re.search(title_pattern, log.Session_Title, re.I)):
                    last_match = (log, st)
            if last_match:
                log, st = last_match
                if info_key:
                    prefix = info_key.replace("{{", "").replace("_info}}", "")
                    replacements[info_key] = info_fmt(log.owner.Name, derive_credentials(log.owner), prefix.replace("-", " ").title())
                    if avatar_map is not None: avatar_map[f"{prefix}_avatar"] = log.owner
                if dur_key: replacements[dur_key] = dur_fmt(log)
                if info_key:
                    prefix = info_key.replace("{{", "").replace("_info}}", "")
                    replacements[f"{{{{{prefix}_session_title}}}}"] = log.Session_Title or ""
                return True
            return False
        for prefix, role_p, title_p in [("saa", None, "SAA Introduction"), ("president", None, "President's Address"), ("welcome-officer", "Welcome Officer", None), ("tme", "Toastmaster", None), ("timer", "Timer", None), ("ah-counter", "Ah-Counter", None), ("grammarian", "Grammarian", None), ("topicsmaster", "Topicsmaster", None), ("ge", "General Evaluator", None), ("photographer", "Photographer", None)]:
            populate_role(role_p, title_p, "{{" + prefix + "_info}}", "{{" + prefix + "_duration}}")
        excomm = meeting.get_excomm()
        if excomm:
            officers = excomm.get_officers()
            role_to_prefix = {"President": "president", "VPE": "vpe", "VPM": "vpm", "VPPR": "vppr", "Secretary": "secretary", "Treasurer": "treasurer", "SAA": "saa", "Welcome Officer": "welcome-officer"}
            for r, c in officers.items():
                p = role_to_prefix.get(r)
                k = "{{" + p + "_info}}" if p else None
                if p and c and k and not replacements.get(k):
                    replacements[k] = info_fmt(c.Name, derive_credentials(c), p.replace("-", " ").title())
                    if avatar_map is not None: avatar_map[f"{p}_avatar"] = c

    @staticmethod
    def _populate_speakers(context, replacements, info_fmt, dur_fmt, avatar_map=None):
        """Internal v1 helper."""
        speakers = [(log, st) for log, st in context.logs if (st.role and st.role.name == "Prepared Speaker" and log.owner)]
        for i, (log, st) in enumerate(speakers[:6], 1):
            replacements[f"{{{{ps{i}}}}}"] = log.owner.Name
            replacements[f"{{{{ps{i}_info}}}}"] = info_fmt(log.owner.Name, derive_credentials(log.owner), f"Speaker {i}")
            replacements[f"{{{{ps{i}_duration}}}}"] = dur_fmt(log)
            replacements[f"{{{{ps{i}_session_title}}}}"] = log.Session_Title or ""
            if avatar_map is not None: avatar_map[f"ps{i}_avatar"] = log.owner
            details = context.speech_details.get(log.id)
            if details:
                replacements[f"{{{{ps{i}_title}}}}"] = details['speech_title'] or details['project_name'] or ""
                replacements[f"{{{{ps{i}-project_info}}}}"] = f"{details['project_code']} - {details['project_name']}" if details['project_code'] and details['project_name'] else (details['project_name'] or "")
            else: replacements[f"{{{{ps{i}_title}}}}"] = log.Session_Title or ""

    @staticmethod
    def _populate_evaluators(context, replacements, info_fmt, dur_fmt, avatar_map=None):
        """Internal v1 helper."""
        evals = [(log, st) for log, st in context.logs if (st.role and st.role.name == "Individual Evaluator" and log.owner)]
        for i, (log, st) in enumerate(evals[:6], 1):
            replacements[f"{{{{ie{i}_info}}}}"] = info_fmt(log.owner.Name, derive_credentials(log.owner), f"Evaluator {i}")
            replacements[f"{{{{ie{i}_duration}}}}"] = dur_fmt(log)
            replacements[f"{{{{ie{i}_session_title}}}}"] = log.Session_Title or f"Evaluation {i}"
            if avatar_map: avatar_map[f"ie{i}_avatar"] = log.owner

    @staticmethod
    def _populate_featured_session(context, meeting, replacements, info_fmt, dur_fmt, avatar_map=None):
        """Internal v1 helper."""
        featured, m_type = None, (meeting.type or "").strip().lower()
        for l, s in context.logs:
            if m_type and ((s.Title and s.Title.strip().lower() == m_type) or (l.Session_Title and l.Session_Title.strip().lower() == m_type)):
                featured = (l, s); break
        if not featured: featured = next(((l, s) for l, s in context.logs if s.role and s.role.name == "Keynote Speaker"), None)
        if not featured: featured = next(((l, s) for l, s in context.logs if s.role and s.role.name in ["Moderator", "Workshop Presenter", "Moderator-Host"]), None)
        if featured:
            replacements["{{keynote_title}}"] = replacements["{{keynote_session_title}}"] = featured[0].Session_Title or meeting.type or "Featured Session"
            replacements["{{keynote_duration}}"] = dur_fmt(featured[0])
            if featured[0].owner:
                replacements["{{keynote-speaker_info}}"] = info_fmt(featured[0].owner.Name, derive_credentials(featured[0].owner), meeting.type or "Featured Speaker")
                if avatar_map: avatar_map["keynote-speaker_avatar"] = featured[0].owner
        else: replacements["{{keynote_title}}"] = meeting.type or ""

    @staticmethod
    def _populate_table_topics(context, replacements, dur_fmt):
        """Internal v1 helper."""
        tt = next(((l, s) for l, s in context.logs if (s.Title == "Table Topics" or (s.role and s.role.name == "Topicsmaster"))), (None, None))
        if tt[0]:
            replacements["{{table-topics_duration}}"] = dur_fmt(tt[0])
            replacements["{{table-topics_session_title}}"] = tt[0].Session_Title or "Table Topics"

    @staticmethod
    def _identify_host_type(slide):
        """Internal v1 helper."""
        role_map = {"moderator": ["moderator", "panel discussion"], "timer": ["timer", "timing record"], "ah-counter": ["ah-counter", "ah counter"], "grammarian": ["grammarian"], "topicsmaster": ["topicsmaster", "table topics"], "ge": ["general evaluator", " ge "], "saa": ["saa", "sergeant at arms"], "tme": ["tme", "toastmaster"], "president": ["president"], "debate host": ["debate host", "debate"]}
        for shape in slide.shapes:
            name = (shape.name or "").lower()
            for r, kw in role_map.items():
                if any(k in name for k in kw): return r
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            text = shape.text.lower()
            for r, kw in role_map.items():
                if any(k in text for k in kw): return r
        return None

    @staticmethod
    def _populate_slide(slide, meeting, context, replacements, avatar_map, host_type=None):
        """Internal v1 helper."""
        def info_fmt(n, c, r=None):
            if not n: return ""
            fn = f"{n}, {c}" if c else n
            return f"{fn}\n{r}" if r else fn
        def dur_fmt(l):
            if not l: return ""
            dm, dx = l.Duration_Min, l.Duration_Max
            if dm is not None and dx is not None:
                if dm == 0 and dx == 0: return ""
                if dm == 0: return f"{dx} '"
                return f"{dm} ~ {dx} '"
            v = dx if dx is not None else dm
            return f"{v} '" if v and v != 0 else ""
        lr, lam = replacements.copy(), (avatar_map.copy() if avatar_map else {})
        if host_type:
            p = host_type.replace(" ", "-")
            hk = p.replace("-", "_")
            dv = replacements.get(f"{{{{{p}_duration}}}}", "")
            if not dv:
                target = "Moderator" if host_type == "moderator" else ("Debate Host" if host_type == "debate host" else None)
                hl = next((l for l, s in context.logs if s.role and s.role.name == target), None) if target else None
                dv = dur_fmt(hl) if hl else ""
            lr["{{duration}}"] = lr["{{slide_duration}}"] = dv
            iv, st = replacements.get(f"{{{{{hk}_info}}}}", ""), replacements.get(f"{{{{{hk}_session_title}}}}", "")
            if not iv:
                t = "Moderator" if host_type == "moderator" else ("Debate Host" if host_type == "debate host" else None)
                hl = next((l for l, s in context.logs if s.role and s.role.name == t), None) if t else None
                if hl and hl.owner:
                    iv = info_fmt(hl.owner.Name, derive_credentials(hl.owner), t)
                    st = hl.Session_Title or t
                    lam["host_avatar"] = hl.owner
                    lr[f"{t.lower().replace(' ', '_')}_info"] = iv
            lr["{{host_info}}"] = iv
            lr["{{slide_title}}"] = st.upper() if st else host_type.upper()
            for s in slide.shapes:
                if s.is_placeholder and s.placeholder_format.type == 1:
                    if not s.text or "click to add title" in s.text.lower() or s.text.lower() == host_type.lower():
                        s.text = st.upper() if st else host_type.upper()
                        if s.has_text_frame: s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if f"{hk}_avatar" in lam: lam["host_avatar"] = lam[f"{hk}_avatar"]
            er = "Panelist" if host_type == "moderator" else ("Debater" if host_type == "debate host" else None)
            if er:
                ex = []
                for l, s in context.logs:
                    if s.role and s.role.name == er:
                        if l.owners: ex.extend(l.owners)
                        elif l.owner: ex.append(l.owner)
                for i, c in enumerate(ex[:10], 1):
                    info = info_fmt(c.Name, derive_credentials(c), er)
                    lr[f"{{{{contact{i}_info}}}}"] = lr[f"{{{{{er.lower()}{i}_info}}}}"] = info
                    lam[f"contact{i}_avatar"] = c
        cp = set()
        for s in slide.shapes:
            if s.has_text_frame: cp.update(re.findall(r'\{\{.*?_info\}\}', s.text))
        if len(cp) <= 1:
            for k in cp:
                if k in lr and "\n" in lr[k]: lr[k] = lr[k].split("\n", 1)[0]
            if "{{host_info}}" in lr and "{{host_info}}" in cp and "\n" in lr["{{host_info}}"]:
                lr["{{host_info}}"] = lr["{{host_info}}"].split("\n", 1)[0]
        def rr(t):
            t = re.sub(r'Evaluator\s*for', 'Evaluator for', t, flags=re.I)
            t = re.sub(r'INDIVIDUAL EVALUATOR\s*for', 'INDIVIDUAL EVALUATOR for', t, flags=re.I)
            for k, v in lr.items():
                if k in t: t = t.replace(k, str(v))
            return t
        for s in slide.shapes:
            if s.name in lam or not s.has_text_frame: continue
            for p in s.text_frame.paragraphs:
                ft = "".join(r.text for r in p.runs); ut = rr(ft)
                if ut != ft:
                    if p.runs:
                        p.runs[0].text = ut
                        for i in range(1, len(p.runs)): p.runs[i].text = ""
                    else: p.text = ut
        MeetingSlideService._replace_avatar_shapes(slide, lam)

    @staticmethod
    def _crop_image_to_aspect_ratio(image_path, target_width, target_height):
        """Internal v1 helper."""
        try:
            im = Image.open(image_path)
            iw, ih = im.size
            tr, ir = target_width / target_height, iw / ih
            if ir > tr:
                nw = int(ih * tr)
                l = (iw - nw) // 2
                box = (l, 0, l + nw, ih)
            else:
                nh = int(iw / tr)
                t = (ih - nh) // 2
                box = (0, t, iw, t + nh)
            ci = im.crop(box)
            if ci.mode in ('RGBA', 'LA'):
                bg = Image.new(ci.mode[:-1], ci.size, (255, 255, 255))
                bg.paste(ci, mask=ci.split()[-1])
                ci = bg.convert('RGB')
            elif ci.mode != 'RGB': ci = ci.convert('RGB')
            bp, ex = os.path.splitext(image_path)
            cp = f"{bp}_cropped.jpg"
            ci.save(cp, 'JPEG', quality=95)
            return cp
        except Exception as e:
            current_app.logger.error(f"Error cropping image: {e}")
            return None

    @staticmethod
    def _fill_shape_with_image(slide, shape, image_path):
        """Internal v1 helper."""
        try:
            dp = slide.shapes.add_picture(image_path, 0, 0, 0, 0)
            rid = dp._element.blipFill.blip.rEmbed
            slide.shapes._spTree.remove(dp._element)
            sppr = shape._element.spPr
            for c in list(sppr):
                if c.tag.endswith(('solidFill', 'gradFill', 'noFill', 'blipFill', 'pattFill', 'grpFill')): sppr.remove(c)
            nsr = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
            bf, bl, st, fr = OxmlElement('a:blipFill'), OxmlElement('a:blip'), OxmlElement('a:stretch'), OxmlElement('a:fillRect')
            bf.set('rotWithShape', '1'); bl.set(f'{{{nsr}}}embed', rid); st.append(fr); bf.append(bl); bf.append(st)
            idx = 0
            for i, c in enumerate(sppr):
                if c.tag.endswith(('prstGeom', 'custGeom', 'xfrm')): idx = i + 1
            sppr.insert(idx, bf)
        except Exception as e: current_app.logger.error(f"Error filling shape {shape.name}: {e}")

    @staticmethod
    def _replace_avatar_shapes(slide_or_prs, avatar_map):
        """Internal v1 helper."""
        if not avatar_map: return
        nm = {k.strip().lower(): v for k, v in avatar_map.items()}
        if hasattr(slide_or_prs, 'slides') and not hasattr(slide_or_prs, 'shapes'): ss = slide_or_prs.slides
        elif hasattr(slide_or_prs, 'shapes'): ss = [slide_or_prs]
        else: ss = slide_or_prs if isinstance(slide_or_prs, (list, tuple)) else [slide_or_prs]
        for s in ss:
            for sh in s.shapes:
                nc = (sh.name or "").strip().lower()
                if nc in nm:
                    d = nm[nc]
                    url = d.Avatar_URL if d and hasattr(d, 'Avatar_URL') else None
                    ip = None
                    if url:
                        if url.startswith('/static/'): url = url[8:]
                        elif url.startswith('static/'): url = url[7:]
                        if '/' not in url and '\\' not in url:
                            url = os.path.join(current_app.config.get('AVATAR_ROOT_DIR', 'uploads/avatars'), url)
                        ip = os.path.join(current_app.static_folder, url.lstrip('/'))
                        if not os.path.exists(ip): ip = None
                    if not ip:
                        for dn in ["default_avatar.jpg", "default_avatar.png", "avatar_default.jpg"]:
                            tp = os.path.join(current_app.static_folder, dn)
                            if os.path.exists(tp): ip = tp; break
                        if not ip:
                            tp = os.path.join(current_app.static_folder, current_app.config.get('AVATAR_ROOT_DIR', 'uploads/avatars'), "default.jpg")
                            if os.path.exists(tp): ip = tp
                    try:
                        cp = MeetingSlideService._crop_image_to_aspect_ratio(ip, sh.width, sh.height)
                        if cp:
                            MeetingSlideService._fill_shape_with_image(s, sh, cp)
                            try: os.remove(cp)
                            except: pass
                    except Exception as e: current_app.logger.error(f"Error processing avatar for {sh.name}: {e}")
