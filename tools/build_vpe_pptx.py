"""Generate the comprehensive VPE tutorial PPTX.

Usage:  python tools/build_vpe_pptx.py
Reads screenshots from /tmp/mm_snapshots, writes to docs/VPE_MemMaker_Tutorial.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------
SNAP = "/Users/wmu/workspace/toastmasters/vpemaster/docs/assets/vpe_tutorial_snapshots"
OUT  = "/Users/wmu/workspace/toastmasters/vpemaster/docs/VPE_MemMaker_Tutorial.pptx"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
MAROON       = RGBColor(0x8B, 0x1A, 0x1A)
DARK_MAROON  = RGBColor(0x5A, 0x11, 0x11)
GOLD         = RGBColor(0xD4, 0xA0, 0x4C)
DARK_TEXT    = RGBColor(0x2C, 0x2C, 0x2C)
GRAY_TEXT    = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_BG     = RGBColor(0xF7, 0xF3, 0xEE)   # warm off-white
PANEL_BG     = RGBColor(0xFA, 0xFA, 0xFA)
BORDER       = RGBColor(0xD8, 0xD8, 0xD8)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GREEN        = RGBColor(0x2E, 0x7D, 0x32)
RED          = RGBColor(0xC6, 0x28, 0x28)
BLUE         = RGBColor(0x1E, 0x40, 0x7A)

FONT = "Calibri"
FONT_TITLE = "Calibri"

# ---------------------------------------------------------------------------
# Presentation setup
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # blank layout — full control


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w is not None:
            s.line.width = line_w
    return s


def add_text(slide, text, x, y, w, h, *, size=18, bold=False, color=DARK_TEXT,
             font=FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    if isinstance(text, str):
        lines = [text]
    else:
        lines = list(text)

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullets(slide, items, x, y, w, h, *, size=18, color=DARK_TEXT,
                bullet_color=MAROON, spacing=6, line_spacing=1.15, indent_em=0.18):
    """items: list of strings or (label, body) tuples for bold lead-in style."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(spacing)
        # bullet marker
        r0 = p.add_run()
        r0.text = "■  "
        r0.font.name = FONT
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            label, body = item
            r1 = p.add_run()
            r1.text = label
            r1.font.name = FONT
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = DARK_TEXT
            if body:
                r2 = p.add_run()
                r2.text = " — " + body
                r2.font.name = FONT
                r2.font.size = Pt(size)
                r2.font.color.rgb = DARK_TEXT
        else:
            r1 = p.add_run()
            r1.text = item
            r1.font.name = FONT
            r1.font.size = Pt(size)
            r1.font.color.rgb = DARK_TEXT
    return tb


def add_image(slide, path, x, y, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(path, x, y, width=w, height=h)
    if w:
        return slide.shapes.add_picture(path, x, y, width=w)
    if h:
        return slide.shapes.add_picture(path, x, y, height=h)
    return slide.shapes.add_picture(path, x, y)


def add_image_with_caption(slide, path, caption, x, y, w, h, caption_h=Inches(0.35)):
    """Place a screenshot inside a bordered panel with caption underneath."""
    # panel background
    panel = add_rect(slide, x, y, w, h + caption_h + Inches(0.08),
                     fill=WHITE, line=BORDER, line_w=Pt(0.75))
    # image inside
    add_image(slide, path, x + Inches(0.05), y + Inches(0.05), w=w - Inches(0.10))
    # caption
    add_text(slide, caption, x, y + h + Inches(0.04), w, caption_h,
             size=11, color=GRAY_TEXT, align=PP_ALIGN.CENTER, italic=True)


# ---------------------------------------------------------------------------
# Layout primitives
# ---------------------------------------------------------------------------
def add_background(slide, color=WHITE):
    bg = add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=color)
    bg.shadow.inherit = False
    return bg


def add_corner_accent(slide):
    """A subtle maroon strip at the top + small gold square for brand."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), fill=MAROON)
    add_rect(slide, Inches(0.6), Inches(0.45), Inches(0.18), Inches(0.18), fill=GOLD)


def add_footer(slide, page_num, total):
    # bottom rule
    add_rect(slide, Inches(0.6), Inches(7.15), Inches(12.13), Emu(9525),
             fill=BORDER)
    add_text(slide, "Memory Maker — VPE Tutorial", Inches(0.6), Inches(7.20),
             Inches(6), Inches(0.25), size=10, color=GRAY_TEXT, align=PP_ALIGN.LEFT)
    add_text(slide, f"{page_num} / {total}", Inches(11.5), Inches(7.20),
             Inches(1.23), Inches(0.25), size=10, color=GRAY_TEXT, align=PP_ALIGN.RIGHT)


def add_slide_title(slide, eyebrow, title, subtitle=None):
    add_corner_accent(slide)
    if eyebrow:
        add_text(slide, eyebrow, Inches(0.6), Inches(0.42), Inches(10), Inches(0.3),
                 size=11, bold=True, color=MAROON, align=PP_ALIGN.LEFT)
    add_text(slide, title, Inches(0.6), Inches(0.75), Inches(12.13), Inches(0.7),
             size=32, bold=True, color=DARK_MAROON, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.6), Inches(1.45), Inches(12.13), Inches(0.4),
                 size=14, color=GRAY_TEXT, italic=True, align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
TOTAL = 28   # will be set after build
PG = [0]
def new_slide():
    PG[0] += 1
    return prs.slides.add_slide(BLANK)


def slide_cover():
    s = new_slide()
    # Background block
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=LIGHT_BG)
    # Left maroon band
    add_rect(s, Inches(0), Inches(0), Inches(0.5), SLIDE_H, fill=MAROON)
    # Gold accent square
    add_rect(s, Inches(1.5), Inches(1.4), Inches(0.6), Inches(0.6), fill=GOLD)
    # Eyebrow
    add_text(s, "VPE TUTORIAL", Inches(1.5), Inches(2.15), Inches(8), Inches(0.4),
             size=18, bold=True, color=MAROON)
    # Title
    add_text(s, "Memory Maker", Inches(1.5), Inches(2.6), Inches(11), Inches(1.2),
             size=60, bold=True, color=DARK_MAROON)
    add_text(s, "Managing a Full Meeting Cycle",
             Inches(1.5), Inches(3.7), Inches(11), Inches(0.8),
             size=32, color=DARK_TEXT)
    # Subtitle
    add_text(s, "A step-by-step walkthrough for the Vice President Education",
             Inches(1.5), Inches(4.6), Inches(11), Inches(0.5),
             size=18, color=GRAY_TEXT, italic=True)
    # Decorative line
    add_rect(s, Inches(1.5), Inches(5.4), Inches(2.5), Emu(28575), fill=GOLD)
    # Footer
    add_text(s, "Generated 2026-06-22  ·  Memory Maker BETA  ·  For internal club training",
             Inches(1.5), Inches(6.7), Inches(11), Inches(0.4),
             size=12, color=GRAY_TEXT)


def slide_toc():
    s = new_slide()
    add_slide_title(s, "CONTENTS", "What this deck covers")
    items = [
        ("Phase 1 — Pre-Meeting Planning",  "Slides 4–11  ·  Create the meeting, build the agenda, open booking"),
        ("Phase 2 — Reception",             "Slides 12–13 ·  SAA handles check-in, you stay reachable"),
        ("Phase 3 — Live Meeting",          "Slides 14–17 ·  Publish → Start → react → Stop"),
        ("Phase 4 — Wrap-Up & Records",     "Slides 18–22 ·  Voting results, speech logs, Pathways, archive"),
        ("Reference",                       "Slides 23–28 ·  Status flow, permissions, pitfalls, glossary"),
    ]
    y = 2.0
    for i, (head, sub) in enumerate(items):
        # number circle
        n = add_rect(s, Inches(0.7), Inches(y), Inches(0.55), Inches(0.55),
                     fill=MAROON, shape=MSO_SHAPE.OVAL)
        add_text(s, str(i + 1), Inches(0.7), Inches(y + Inches(0.07)),
                 Inches(0.55), Inches(0.5),
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, head, Inches(1.5), Inches(y), Inches(11), Inches(0.4),
                 size=20, bold=True, color=DARK_TEXT)
        add_text(s, sub, Inches(1.5), Inches(y + Inches(0.4)), Inches(11), Inches(0.4),
                 size=13, color=GRAY_TEXT, italic=True)
        y += 0.9
    add_footer(s, PG[0], TOTAL)


def slide_phase_overview():
    s = new_slide()
    add_slide_title(s, "OVERVIEW", "The four phases of a meeting cycle")
    # 4 cards
    phases = [
        ("1", "Pre-Meeting\nPlanning",   "T-7 → T-1 day",  "Build the agenda\n& open booking", MAROON),
        ("2", "Reception",                "T-0  18:30",    "SAA checks in\nguests & members",   GOLD),
        ("3", "Live Meeting",             "T-0  19:00–21:00","Publish → Start\n→ react → Stop",   DARK_MAROON),
        ("4", "Wrap-Up\n& Records",       "T+0 → T+1",     "Tally votes, log\nspeeches, archive", BLUE),
    ]
    card_w = Inches(2.85)
    gap = Inches(0.18)
    x0 = Inches(0.6)
    y0 = Inches(2.0)
    h  = Inches(4.2)
    for i, (num, name, when, what, color) in enumerate(phases):
        x = x0 + (card_w + gap) * i
        # card body
        add_rect(s, x, y0, card_w, h, fill=PANEL_BG, line=BORDER, line_w=Pt(0.75))
        # color bar at top
        add_rect(s, x, y0, card_w, Inches(0.6), fill=color)
        # big number
        add_text(s, num, x, y0 + Inches(0.05), card_w, Inches(0.5),
                 size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # name
        add_text(s, name, x, y0 + Inches(0.8), card_w, Inches(0.9),
                 size=20, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # when
        add_text(s, when, x, y0 + Inches(2.0), card_w, Inches(0.4),
                 size=12, color=color, bold=True, align=PP_ALIGN.CENTER)
        # what
        add_text(s, what, x, y0 + Inches(2.5), card_w, Inches(1.4),
                 size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    # bottom takeaway
    add_rect(s, Inches(0.6), Inches(6.5), Inches(12.13), Inches(0.55),
             fill=LIGHT_BG, line=BORDER)
    add_text(s, "→  The VPE owns Phase 1 & Phase 3, supports Phase 4, stays reachable in Phase 2.",
             Inches(0.6), Inches(6.55), Inches(12.13), Inches(0.45),
             size=14, bold=True, color=DARK_MAROON, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, PG[0], TOTAL)


# Generic section divider
def slide_section_divider(num, name, subtitle):
    s = new_slide()
    add_background(s, DARK_MAROON)
    # gold accent block
    add_rect(s, Inches(0.6), Inches(2.6), Inches(0.6), Inches(0.6), fill=GOLD)
    add_text(s, f"PHASE {num}", Inches(0.6), Inches(3.3), Inches(12), Inches(0.5),
             size=20, bold=True, color=GOLD)
    add_text(s, name, Inches(0.6), Inches(3.85), Inches(12), Inches(1.2),
             size=54, bold=True, color=WHITE)
    add_rect(s, Inches(0.6), Inches(5.2), Inches(2), Emu(28575), fill=GOLD)
    add_text(s, subtitle, Inches(0.6), Inches(5.4), Inches(12), Inches(0.8),
             size=18, color=RGBColor(0xF0, 0xE0, 0xC0), italic=True)
    add_text(s, f"{PG[0]} / {TOTAL}", Inches(11.5), Inches(7.2),
             Inches(1.23), Inches(0.25), size=10,
             color=RGBColor(0xC0, 0xB0, 0x90), align=PP_ALIGN.RIGHT)


def slide_content(step, title, lead, bullets, image=None, image_caption=None):
    """Standard 2-column content slide: bullets on left, screenshot on right."""
    s = new_slide()
    add_slide_title(s, f"STEP {step}", title, lead)
    # bullets
    add_bullets(s, bullets, Inches(0.6), Inches(2.1), Inches(6.2), Inches(4.6),
                size=15, spacing=10, line_spacing=1.2)
    # image
    if image:
        # Right column: image panel
        add_image_with_caption(s, image, image_caption or "",
                               Inches(7.0), Inches(2.0),
                               Inches(5.8), Inches(4.2))
    add_footer(s, PG[0], TOTAL)


def slide_text_only(step, title, lead, bullets):
    s = new_slide()
    add_slide_title(s, f"STEP {step}", title, lead)
    add_bullets(s, bullets, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5),
                size=16, spacing=10, line_spacing=1.2)
    add_footer(s, PG[0], TOTAL)


# ---------------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------------
slide_cover()
slide_toc()
slide_phase_overview()

# ============ PHASE 1 ============
slide_section_divider(1, "Pre-Meeting Planning",
                      "T-7 to T-1 day  ·  The VPE's main work block")

slide_content(
    "1.1", "Open the Agenda page",
    "Memory Maker's home base for everything meeting-related.",
    [
        ("Route", "/agenda — also reachable from the top nav under Meetings"),
        ("Meeting selector", "Use the dropdown to pick the meeting you want to work on"),
        ("Status group", "Shows the current meeting's status: Unpublished / Not Started / Running / Finished / Cancelled"),
        ("Agenda table", "Lists every session row with time, role, owner, project, and validation flags"),
    ],
    image=os.path.join(SNAP, "agenda_main.png"),
    image_caption="The /agenda page — meeting #982 (Aug 04, 2026), status: Not Started"
)

slide_content(
    "1.2", "Create a meeting from a template",
    "Start a new meeting by filling in the date, theme, and template.",
    [
        ("Click New", "Top of the agenda editor, or via the top-right menu"),
        ("Fill the form", "Date, meeting number, title/theme, subtitle, type, media URL, template"),
        ("Template file", "Determines the default set of agenda items (Standard Club Meeting, etc.)"),
        ("Save", "POST /agenda/create — requires MEETING_CREATE permission"),
    ],
    image=os.path.join(SNAP, "settings_main.png"),
    image_caption="Settings → Modules — make sure 'Data/Slides Export' is on before generating slides"
)

slide_content(
    "1.3", "Build the agenda",
    "Edit sessions inline: time, role, project, and meeting-level fields.",
    [
        ("Add/remove a session", "Use the row controls at the bottom of the table"),
        ("Reorder", "Drag the row handle, or use the up/down arrows"),
        ("Edit start time", "Click the time cell, type 24h HH:MM — overlaps are flagged"),
        ("Word of the Day", "Set in the top-of-page wod-display block"),
        ("GE Style", "The ge_mode dropdown controls the General Evaluator's report format"),
        ("Red badges", "Indicate validation problems — resolve before publishing"),
    ],
    image=os.path.join(SNAP, "agenda_main.png"),
    image_caption="The agenda table — session rows are editable inline"
)

slide_content(
    "1.4", "Add project info for prepared speeches",
    "Prepared speech slots are flagged Valid_for_Project = true and need a Pathways project.",
    [
        ("Click the project cell", "In the row of the prepared-speech session"),
        ("Pick the pathway", "e.g. Dynamic Leadership, Presentation Mastery"),
        ("Pick the project", "e.g. Level 1 — Ice Breaker"),
        ("Add an optional title", "Used on slides and in speech records"),
        ("If undecided", "Leave blank — update before the meeting is Started"),
    ],
    image=os.path.join(SNAP, "pathway_library.png"),
    image_caption="Pathways Library — the catalog of projects you assign to speech slots"
)

slide_content(
    "1.5", "Open role booking",
    "Publishing unlocks role booking for all members.",
    [
        ("Click Publish", "In the agenda page's status group"),
        ("Confirm the prompt", "Memory Maker reminds you that booking opens to all members"),
        ("Status flips", "Unpublished → Not Started"),
        ("Members see /booking", "Self-booking and waitlists become available"),
        ("Tip", "Once published, structural changes to sessions are harder — finalise first"),
    ],
    image=os.path.join(SNAP, "agenda_main.png"),
    image_caption="Status group on the agenda — click Publish to open booking"
)

slide_content(
    "1.6", "Book or assign roles",
    "Members self-book; you backfill empty slots and assign evaluators.",
    [
        ("Member self-booking", "Members visit /booking, pick a meeting, click a slot"),
        ("Waitlists", "Full slots show 'Join Waitlist' — auto-promote when owner cancels"),
        ("Officer assignment", "Click the owner cell on the agenda page and pick a contact"),
        ("Backfill 24–48h before", "Fill any empty required roles; assign evaluators to speeches"),
        ("Tip", "Read the Planner page 1–2 days before — it shows members' intentions"),
    ],
    image=os.path.join(SNAP, "booking.png"),
    image_caption="/booking — the Booking Assistant as members see it"
)

slide_text_only(
    "1.7", "Day-before checklist",
    "Run through this list the day before the meeting.",
    [
        "Agenda has no red validation badges",
        "Every prepared speech has a Pathways project assigned",
        "Every required role has an owner (or is intentionally left empty for the TME)",
        "Word of the Day is set; the WOD speaker knows the definition",
        "Slides have been generated and saved to the laptop",
        "The SAA has confirmed who will publish and handle voting on the night",
        "If all green, you are ready for Phase 3",
    ]
)

# ============ PHASE 2 ============
slide_section_divider(2, "Reception",
                      "T-0  18:30  ·  Owned by the SAA — your job is to stay reachable")

slide_content(
    "2.1", "Reception & check-in",
    "The SAA runs the door; you confirm the meeting is ready to publish.",
    [
        ("Route", "/roster — SAA adds guests and members as they arrive"),
        ("Ticket types", "Early-bird / Walk-in / Officer — order number auto-calculates"),
        ("Live KPI", "Total attendees and revenue shown at the top in real time"),
        ("Your only task", "Stay reachable. Don't publish/start the meeting until the SAA is ready"),
        ("Last-minute edits", "Adding a table topic speaker or roster entry mid-meeting is supported"),
    ],
    image=os.path.join(SNAP, "roster.png"),
    image_caption="/roster — the SAA's check-in tool"
)

# ============ PHASE 3 ============
slide_section_divider(3, "Live Meeting",
                      "T-0  19:00–21:00  ·  The status button is your only control surface")

slide_content(
    "3.1", "Publish and start the meeting",
    "Two clicks: Publish (unpublished → not started), then Start (not started → running).",
    [
        ("Publish", "Click in the status group; status flips to Not Started"),
        ("Wait for 6:55 PM", "Start the meeting 5 minutes before call-to-order"),
        ("Click Start", "Sets Start_Time to now; status flips to Running"),
        ("Voting enables", "Members can submit ballots at /voting; SAA watches live count"),
        ("Reactions enabled", "You can edit owners, add sessions, and add table topic speakers on the fly"),
    ],
    image=os.path.join(SNAP, "agenda_main.png"),
    image_caption="Status group: Publish → Start → Stop"
)

slide_text_only(
    "3.2", "React to in-meeting changes",
    "Common requests and how to handle each in under a minute.",
    [
        ("Add a table topic speaker", "Click the owner cell of the TT session and pick the person"),
        ("Swap a speaker", "Click the owner cell, clear the old, pick the new — no page reload needed"),
        ("Add a TT session on the fly", "Use the Add Session control; drag to reorder if needed"),
        ("Forgot an evaluator", "Same as owner-edit flow — just pick from the contact list"),
        ("All edits via", "POST /agenda/update — one round trip, validated and persisted atomically"),
    ]
)

slide_content(
    "3.3", "End the meeting",
    "Click Stop. Memory Maker seals the meeting and runs the post-meeting bookkeeping.",
    [
        ("Click Stop", "Status flips to Finished"),
        ("Vote tally", "Winners are written to the voting records"),
        ("Waitlist cleanup", "All waitlist entries for this meeting are deleted"),
        ("Planner sync", "booked → completed,  waitlist → obsolete,  draft stays draft"),
        ("Project auto-complete", "Every project's session log is marked Completed and metadata synced"),
    ],
    image=os.path.join(SNAP, "voting.png"),
    image_caption="/voting — auto-tallied when the meeting stops"
)

# ============ PHASE 4 ============
slide_section_divider(4, "Wrap-Up & Records",
                      "T+0 to T+1  ·  Confirm votes, log speeches, archive the meeting")

slide_content(
    "4.1", "Confirm voting results",
    "The SAA reads final results and runs the lucky draw; you verify the data.",
    [
        ("Route", "/voting — final tallies and award winners"),
        ("NPS", "/voting_nps — club-level satisfaction trend"),
        ("SAA's verification list", "Confirms the right people got the right prizes"),
        ("If something looks wrong", "Don't edit ballots directly — file an issue from /issues to preserve the audit trail"),
    ],
    image=os.path.join(SNAP, "voting.png"),
    image_caption="/voting — final tallies after Stop"
)

slide_content(
    "4.2", "Update the media URL",
    "Attach the meeting recording so members can find it later.",
    [
        ("Back on /agenda", "Edit the meeting's Media URL field in the header"),
        ("Save", "The recording is now linked to the meeting record"),
        ("Where it shows up", "Members see it on their Planner and in the meeting archive"),
    ],
    image=os.path.join(SNAP, "contacts.png"),
    image_caption="/contacts — review each speaker's record"
)

slide_content(
    "4.3", "Review speech logs and Pathways",
    "Auto-completion usually does the work; verify a few key cases.",
    [
        ("/contacts", "Confirm each speaker's project and level are recorded correctly"),
        ("Auto-completion", "Project Status flipped to Completed when the meeting was stopped"),
        ("If a project looks wrong", "Fix it via /speech_logs and re-save — metadata re-syncs"),
        ("Pathways progress", "Visit /pathway_library to confirm a path is selected; level progress increments automatically"),
        ("Common bug", "Generic placeholder project (ProjectID.GENERIC) blocks the level increment — replace it"),
    ],
    image=os.path.join(SNAP, "pathway_library.png"),
    image_caption="Pathways Library — the catalog that backs the level tracking"
)

slide_text_only(
    "4.4", "Achievements, archive, and planning",
    "Optional wrap-up steps and the planning handoff to the next cycle.",
    [
        ("Achievements", "If your club uses /achievements, award any new badges (e.g. First Ice Breaker)"),
        ("Archive or delete", "Leave finished meetings in the archive (default), or click Delete to hard-remove"),
        ("Delete requires", "MEETING_CREATE permission; no undo — waitlist, votes, planner entries all go with it"),
        ("Plan ahead", "Update templates and settings if this meeting set a new precedent (see template-manager-plan.md)"),
        ("Brief the next VPE", "Pass along any quirks (custom sessions, recurring issues) before term handover"),
    ]
)

# ============ REFERENCE ============
def slide_status_flow():
    s = new_slide()
    add_slide_title(s, "REFERENCE", "Meeting status flow",
                    "One-way transitions. To re-do a meeting, create a new one.")

    # Status nodes
    nodes = [
        ("unpublished", MAROON,  "Unpublished"),
        ("not started", GOLD,    "Not Started"),
        ("running",     DARK_MAROON, "Running"),
        ("finished",    BLUE,    "Finished"),
    ]
    # Edges
    edges = [
        ("Publish",   0, 1),
        ("Start",     1, 2),
        ("Stop",      2, 3),
    ]

    # Layout: 4 nodes evenly spaced, vertically centered
    y = Inches(3.2)
    node_w = Inches(2.4)
    node_h = Inches(1.1)
    gap = (SLIDE_W - node_w * 4) / 5
    x_positions = [gap + i * (node_w + gap) for i in range(4)]

    for i, (key, color, label) in enumerate(nodes):
        x = x_positions[i]
        add_rect(s, x, y, node_w, node_h, fill=color, line=color)
        add_text(s, label, x, y, node_w, node_h,
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Arrows + edge labels
    for label, src, dst in edges:
        x1 = x_positions[src] + node_w
        x2 = x_positions[dst]
        cy = y + node_h / 2
        # arrow line
        line = s.shapes.add_connector(1, x1, cy, x2 - Inches(0.3), cy)  # straight
        line.line.color.rgb = GRAY_TEXT
        line.line.width = Pt(1.5)
        # arrow head
        ln = line.line._get_or_add_ln()
        tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
        tailEnd.set('type', 'triangle')
        # label above
        add_text(s, label, x1, y - Inches(0.55), x2 - x1, Inches(0.4),
                 size=14, bold=True, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    # Delete note
    add_rect(s, Inches(0.6), Inches(5.5), Inches(12.13), Inches(0.9),
             fill=PANEL_BG, line=BORDER)
    add_text(s,
             "From Finished → Delete (only if MEETING_CREATE permission).  Hard delete: votes, roster, "
             "waitlist, planner entries, and session logs all removed. No undo.",
             Inches(0.8), Inches(5.55), Inches(11.93), Inches(0.8),
             size=13, color=DARK_TEXT, italic=True, anchor=MSO_ANCHOR.MIDDLE)

    # Side-effects summary below
    add_text(s, "Auto-actions on Stop:",
             Inches(0.6), Inches(6.5), Inches(12.13), Inches(0.3),
             size=13, bold=True, color=MAROON)
    add_text(s,
             "votes tallied  ·  waitlist cleared  ·  Planner booked→completed  ·  "
             "projects auto-completed  ·  contact metadata re-synced",
             Inches(0.6), Inches(6.78), Inches(12.13), Inches(0.3),
             size=12, color=GRAY_TEXT, italic=True)
    add_footer(s, PG[0], TOTAL)


def slide_permissions():
    s = new_slide()
    add_slide_title(s, "REFERENCE", "Permissions cheat sheet",
                    "Most of this deck assumes the standard VPE role grants these.")

    headers = ["Action", "Permission"]
    rows = [
        ("Create a meeting",                          "MEETING_CREATE"),
        ("Edit agenda (sessions, WOD, projects)",     "MEETING_MANAGE"),
        ("Assign / remove role owners",               "MEETING_MANAGE"),
        ("Publish / Start / Stop the meeting",        "MEETING_MANAGE"),
        ("Delete a finished meeting",                 "MEETING_CREATE"),
        ("View all meetings (incl. unpublished)",     "MEETING_VIEW_ALL"),
        ("Generate slides / export",                  "MEETING_MANAGE + module 'Data/Slides Export'"),
    ]

    # Table
    cols = [Inches(7.0), Inches(5.13)]
    row_h = Inches(0.45)
    x0 = Inches(0.6)
    y0 = Inches(2.1)
    # Header
    add_rect(s, x0, y0, cols[0] + cols[1], row_h, fill=DARK_MAROON)
    cx = x0
    for i, h in enumerate(headers):
        add_text(s, h, cx + Inches(0.15), y0, cols[i], row_h,
                 size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        cx += cols[i]
    # Body
    for r, (a, p) in enumerate(rows):
        ry = y0 + row_h * (r + 1)
        fill = PANEL_BG if r % 2 == 0 else WHITE
        add_rect(s, x0, ry, cols[0] + cols[1], row_h, fill=fill, line=BORDER, line_w=Pt(0.5))
        cx = x0
        for i, val in enumerate([a, p]):
            color = DARK_TEXT if i == 0 else MAROON
            bold = i == 1
            add_text(s, val, cx + Inches(0.15), ry, cols[i], row_h,
                     size=13, color=color, bold=bold, anchor=MSO_ANCHOR.MIDDLE)
            cx += cols[i]
    # Footer note
    add_text(s,
             "If a button is greyed out, you are missing the corresponding permission. "
             "Ask a ClubAdmin or check Settings → Users.",
             Inches(0.6), Inches(6.4), Inches(12.13), Inches(0.5),
             size=12, color=GRAY_TEXT, italic=True)
    add_footer(s, PG[0], TOTAL)


def slide_pitfalls():
    s = new_slide()
    add_slide_title(s, "REFERENCE", "Common pitfalls",
                    "Things that have bitten VPEs before. Avoid these.")

    items = [
        ("Publishing too early",
         "Members can book roles as soon as the meeting is published. Finalise structure first."),
        ("Wrong Pathways project",
         "A wrong project pollutes the speaker's record. Ask if unsure — don't guess."),
        ("Clicking Stop by accident",
         "The meeting is sealed. The only recovery path is cloning the meeting or restoring from backup."),
        ("Forgetting to assign an evaluator",
         "Backfill the night before. Evaluators are just owners of evaluator-type sessions."),
        ("Editing the date after publishing",
         "Meeting Number auto-renames; the poster file is renamed too. Verify both after the change."),
        ("Generic project on a prepared speech",
         "ProjectID.GENERIC blocks the level increment. Replace with the real project and re-save."),
    ]
    add_bullets(s, items, Inches(0.7), Inches(2.1), Inches(11.93), Inches(4.6),
                size=15, spacing=10, line_spacing=1.2)
    add_footer(s, PG[0], TOTAL)


def slide_glossary():
    s = new_slide()
    add_slide_title(s, "REFERENCE", "Glossary",
                    "Terms used throughout this deck.")

    items = [
        ("VPE",   "Vice President Education — owns the meeting's content and Pathways tracking."),
        ("SAA",   "Sergeant at Arms — owns the room, check-in, voting tallying, and lucky draw."),
        ("TME",   "Toastmaster of the Evening — runs the meeting on the night."),
        ("WOD",   "Word of the Day — the Memory Maker's traditional responsibility; used by speakers and TME."),
        ("Pathway","Toastmasters' learning experience, made up of projects across multiple levels."),
        ("Session Log", "A row in the agenda table — one per agenda item (TME, Table Topics, Speech, Eval, etc.)."),
        ("Planner", "A member's personal plan: which project they intend to deliver in which meeting."),
        ("Waitlist", "Auto-managed queue for full slots. Next person is promoted when the owner cancels."),
    ]
    add_bullets(s, items, Inches(0.7), Inches(2.1), Inches(11.93), Inches(4.6),
                size=15, spacing=8, line_spacing=1.15)
    add_footer(s, PG[0], TOTAL)


def slide_where_next():
    s = new_slide()
    add_slide_title(s, "REFERENCE", "Where to go next",
                    "Related docs in this repository.")

    rows = [
        ("SAA — checklist view",   "docs/SAA_MemMaker_User_Manual.md"),
        ("SAA — grid view",        "docs/SAA_User_Manual_Grid.md"),
        ("All officers at-a-glance","docs/getting_started.md"),
        ("Pathways deep-dive",     "docs/ (search 'pathways')"),
        ("Template editing",       "docs/template-manager-plan.md"),
        ("Voting / live updates",  "docs/MAILBOX_REALTIME_DESIGN.md"),
        ("User / contact model",   "docs/CONTACT_USER_CLUB_MODEL.md"),
        ("This deck's source doc", "docs/VPE_MemMaker_Tutorial.md"),
    ]
    y = Inches(2.1)
    row_h = Inches(0.42)
    for i, (label, path) in enumerate(rows):
        ry = y + row_h * i
        fill = PANEL_BG if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.6), ry, Inches(12.13), row_h, fill=fill, line=BORDER, line_w=Pt(0.4))
        add_text(s, label, Inches(0.8), ry, Inches(5.5), row_h,
                 size=14, bold=True, color=DARK_TEXT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, path, Inches(6.4), ry, Inches(6.3), row_h,
                 size=13, color=MAROON, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, PG[0], TOTAL)


def slide_closing():
    s = new_slide()
    add_background(s, DARK_MAROON)
    add_rect(s, Inches(0), Inches(0), Inches(0.5), SLIDE_H, fill=GOLD)
    add_rect(s, Inches(1.5), Inches(2.7), Inches(0.6), Inches(0.6), fill=GOLD)
    add_text(s, "YOU'RE READY", Inches(1.5), Inches(3.4), Inches(11), Inches(0.6),
             size=24, bold=True, color=GOLD)
    add_text(s, "Manage the cycle with confidence.",
             Inches(1.5), Inches(4.0), Inches(11), Inches(1.0),
             size=40, bold=True, color=WHITE)
    add_text(s,
             "Plan with care, publish once, stay calm on the night, and "
             "let Memory Maker do the bookkeeping at the end.",
             Inches(1.5), Inches(5.0), Inches(11), Inches(1.0),
             size=18, color=RGBColor(0xF0, 0xE0, 0xC0), italic=True)
    add_rect(s, Inches(1.5), Inches(6.0), Inches(2), Emu(28575), fill=GOLD)
    add_text(s, "Questions? Ping a ClubAdmin or open an issue from /issues.",
             Inches(1.5), Inches(6.2), Inches(11), Inches(0.5),
             size=14, color=RGBColor(0xD0, 0xC0, 0xA0), italic=True)
    add_text(s, f"{PG[0]} / {TOTAL}", Inches(11.5), Inches(7.2),
             Inches(1.23), Inches(0.25), size=10,
             color=RGBColor(0xC0, 0xB0, 0x90), align=PP_ALIGN.RIGHT)


# Build the rest of the deck in order
slide_status_flow()
slide_permissions()
slide_pitfalls()
slide_glossary()
slide_where_next()
slide_closing()

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
prs.save(OUT)
print(f"Wrote {OUT}  ({PG[0]} slides, {os.path.getsize(OUT)//1024} KB)")
