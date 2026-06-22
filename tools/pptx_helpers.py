"""Shared layout helpers for the VPE tutorial PPTX builders.

Each per-language script (build_vpe_pptx.py, build_vpe_pptx_zh.py) imports these
helpers and supplies the localized content.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
MAROON       = RGBColor(0x8B, 0x1A, 0x1A)
DARK_MAROON  = RGBColor(0x5A, 0x11, 0x11)
GOLD         = RGBColor(0xD4, 0xA0, 0x4C)
DARK_TEXT    = RGBColor(0x2C, 0x2C, 0x2C)
GRAY_TEXT    = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_BG     = RGBColor(0xF7, 0xF3, 0xEE)
PANEL_BG     = RGBColor(0xFA, 0xFA, 0xFA)
BORDER       = RGBColor(0xD8, 0xD8, 0xD8)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GREEN        = RGBColor(0x2E, 0x7D, 0x32)
RED          = RGBColor(0xC6, 0x28, 0x28)
BLUE         = RGBColor(0x1E, 0x40, 0x7A)
WARM_LIGHT   = RGBColor(0xF0, 0xE0, 0xC0)
WARM_MID     = RGBColor(0xD0, 0xC0, 0xA0)
WARM_DIM     = RGBColor(0xC0, 0xB0, 0x90)

FONT = "Calibri"


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w is not None:
            s.line.width = line_w
    return s


def add_text(slide, text, x, y, w, h, *, size=18, bold=False, color=DARK_TEXT,
             font=FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = [text] if isinstance(text, str) else list(text)

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullets(slide, items, x, y, w, h, *, size=18, color=DARK_TEXT,
                bullet_color=MAROON, spacing=6, line_spacing=1.15):
    """items: list of strings, or (label, body) tuples for bold lead-in style."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(spacing)
        # bullet marker
        r0 = p.add_run()
        r0.text = "■  "
        r0.font.name = FONT
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            label, body = item
            r1 = p.add_run()
            r1.text = label
            r1.font.name = FONT
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = DARK_TEXT
            if body:
                r2 = p.add_run()
                r2.text = " — " + body
                r2.font.name = FONT
                r2.font.size = Pt(size)
                r2.font.color.rgb = DARK_TEXT
        else:
            r1 = p.add_run()
            r1.text = item
            r1.font.name = FONT
            r1.font.size = Pt(size)
            r1.font.color.rgb = DARK_TEXT
    return tb


def add_image(slide, path, x, y, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(path, x, y, width=w, height=h)
    if w:
        return slide.shapes.add_picture(path, x, y, width=w)
    if h:
        return slide.shapes.add_picture(path, x, y, height=h)
    return slide.shapes.add_picture(path, x, y)


def add_image_with_caption(slide, path, caption, x, y, w, h, caption_h=Inches(0.35)):
    """Place a screenshot inside a bordered panel with caption underneath."""
    add_rect(slide, x, y, w, h + caption_h + Inches(0.08),
             fill=WHITE, line=BORDER, line_w=Pt(0.75))
    add_image(slide, path, x + Inches(0.05), y + Inches(0.05), w=w - Inches(0.10))
    add_text(slide, caption, x, y + h + Inches(0.04), w, caption_h,
             size=11, color=GRAY_TEXT, align=PP_ALIGN.CENTER, italic=True)


# ---------------------------------------------------------------------------
# Layout primitives
# ---------------------------------------------------------------------------
def add_background(slide, color=WHITE):
    bg = add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=color)
    bg.shadow.inherit = False
    return bg


def add_corner_accent(slide):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), fill=MAROON)
    add_rect(slide, Inches(0.6), Inches(0.45), Inches(0.18), Inches(0.18), fill=GOLD)


def add_footer(slide, page_num, total):
    add_rect(slide, Inches(0.6), Inches(7.15), Inches(12.13), Emu(9525),
             fill=BORDER)
    add_text(slide, "Memory Maker — VPE Tutorial", Inches(0.6), Inches(7.20),
             Inches(6), Inches(0.25), size=10, color=GRAY_TEXT, align=PP_ALIGN.LEFT)
    add_text(slide, f"{page_num} / {total}", Inches(11.5), Inches(7.20),
             Inches(1.23), Inches(0.25), size=10, color=GRAY_TEXT, align=PP_ALIGN.RIGHT)


def add_slide_title(slide, eyebrow, title, subtitle=None):
    add_corner_accent(slide)
    if eyebrow:
        add_text(slide, eyebrow, Inches(0.6), Inches(0.42), Inches(10), Inches(0.3),
                 size=11, bold=True, color=MAROON, align=PP_ALIGN.LEFT)
    add_text(slide, title, Inches(0.6), Inches(0.75), Inches(12.13), Inches(0.7),
             size=32, bold=True, color=DARK_MAROON, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.6), Inches(1.45), Inches(12.13), Inches(0.4),
                 size=14, color=GRAY_TEXT, italic=True, align=PP_ALIGN.LEFT)


def add_horizontal_arrow(slide, x1, y, x2, color=GRAY_TEXT, width=Pt(1.5)):
    """Draw a horizontal arrow from (x1, y) to (x2, y) with a triangle head."""
    line = slide.shapes.add_connector(1, x1, y, x2 - Inches(0.3), y)
    line.line.color.rgb = color
    line.line.width = width
    ln = line.line._get_or_add_ln()
    tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'triangle')
    return line
