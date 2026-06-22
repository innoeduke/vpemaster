"""Generate the Chinese VPE tutorial PPTX.

Reads screenshots from docs/assets/vpe_tutorial_snapshots,
writes to docs/VPE_MemMaker_Tutorial.zh-CN.pptx
"""
import os
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pptx_helpers import (
    new_presentation, add_rect, add_text, add_bullets, add_image,
    add_image_with_caption, add_background, add_corner_accent, add_footer,
    add_slide_title, add_horizontal_arrow,
    SLIDE_W, SLIDE_H, MAROON, DARK_MAROON, GOLD, DARK_TEXT, GRAY_TEXT,
    LIGHT_BG, PANEL_BG, BORDER, WHITE, BLUE, WARM_LIGHT, WARM_DIM, FONT,
)

SNAP = "/Users/wmu/workspace/toastmasters/vpemaster/docs/assets/vpe_tutorial_snapshots"
OUT  = "/Users/wmu/workspace/toastmasters/vpemaster/docs/VPE_MemMaker_Tutorial.zh-CN.pptx"

prs = new_presentation()
BLANK = prs.slide_layouts[6]

TOTAL = 28
PG = [0]
def new_slide():
    PG[0] += 1
    return prs.slides.add_slide(BLANK)


# ---------------------------------------------------------------------------
# Cover & intro
# ---------------------------------------------------------------------------
def slide_cover():
    s = new_slide()
    add_background(s, LIGHT_BG)
    add_rect(s, Inches(0), Inches(0), Inches(0.5), SLIDE_H, fill=MAROON)
    add_rect(s, Inches(1.5), Inches(1.4), Inches(0.6), Inches(0.6), fill=GOLD)
    add_text(s, "VPE 教程", Inches(1.5), Inches(2.15), Inches(11), Inches(0.4),
             size=18, bold=True, color=MAROON)
    add_text(s, "Memory Maker", Inches(1.5), Inches(2.6), Inches(11), Inches(1.2),
             size=60, bold=True, color=DARK_MAROON)
    add_text(s, "管理一场完整的会议周期",
             Inches(1.5), Inches(3.7), Inches(11), Inches(0.8),
             size=32, color=DARK_TEXT)
    add_text(s, "面向教育副主席 (VPE) 的分步操作指南",
             Inches(1.5), Inches(4.6), Inches(11), Inches(0.5),
             size=18, color=GRAY_TEXT, italic=True)
    add_rect(s, Inches(1.5), Inches(5.4), Inches(2.5), Emu(28575), fill=GOLD)
    add_text(s, "生成于 2026-06-22  ·  Memory Maker BETA  ·  俱乐部内部培训使用",
             Inches(1.5), Inches(6.7), Inches(11), Inches(0.4),
             size=12, color=GRAY_TEXT)


def slide_toc():
    s = new_slide()
    add_slide_title(s, "CONTENTS  ·  目录", "本教程涵盖的内容")
    items = [
        ("阶段一 — 会前规划",   "幻灯片 4–11  ·  创建会议、搭建议程、开放角色预订"),
        ("阶段二 — 接待签到",   "幻灯片 12–13 ·  SAA 负责签到，VPE 保持畅通"),
        ("阶段三 — 现场会议",   "幻灯片 14–17 ·  发布 → 启动 → 应对变更 → 结束"),
        ("阶段四 — 收尾与记录", "幻灯片 18–22 ·  唱票、演讲记录、Pathways、归档"),
        ("参考",                "幻灯片 23–28 ·  状态流转、权限、常见坑、术语表"),
    ]
    y = 2.0
    for i, (head, sub) in enumerate(items):
        n = add_rect(s, Inches(0.7), Inches(y), Inches(0.55), Inches(0.55),
                     fill=MAROON, shape=MSO_SHAPE.OVAL)
        add_text(s, str(i + 1), Inches(0.7), Inches(y + Inches(0.07)),
                 Inches(0.55), Inches(0.5),
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, head, Inches(1.5), Inches(y), Inches(11), Inches(0.4),
                 size=20, bold=True, color=DARK_TEXT)
        add_text(s, sub, Inches(1.5), Inches(y + Inches(0.4)), Inches(11), Inches(0.4),
                 size=13, color=GRAY_TEXT, italic=True)
        y += 0.9
    add_footer(s, PG[0], TOTAL)


def slide_phase_overview():
    s = new_slide()
    add_slide_title(s, "OVERVIEW  ·  概览", "会议周期的四个阶段")
    phases = [
        ("1", "会前规划",  "T-7 → T-1 天",  "搭建议程并\n开放预订",       MAROON),
        ("2", "接待签到",  "T-0  18:30",   "SAA 为宾客和\n会员签到",      GOLD),
        ("3", "现场会议",  "T-0  19:00–21:00","发布 → 启动\n→ 应对 → 结束",DARK_MAROON),
        ("4", "收尾与记录","T+0 → T+1",    "唱票、记录演讲、\n归档",        BLUE),
    ]
    card_w = Inches(2.85)
    gap = Inches(0.18)
    x0 = Inches(0.6)
    y0 = Inches(2.0)
    h  = Inches(4.2)
    for i, (num, name, when, what, color) in enumerate(phases):
        x = x0 + (card_w + gap) * i
        add_rect(s, x, y0, card_w, h, fill=PANEL_BG, line=BORDER, line_w=Pt(0.75))
        add_rect(s, x, y0, card_w, Inches(0.6), fill=color)
        add_text(s, num, x, y0 + Inches(0.05), card_w, Inches(0.5),
                 size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, name, x, y0 + Inches(0.8), card_w, Inches(0.9),
                 size=20, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, when, x, y0 + Inches(2.0), card_w, Inches(0.4),
                 size=12, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, what, x, y0 + Inches(2.5), card_w, Inches(1.4),
                 size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(0.6), Inches(6.5), Inches(12.13), Inches(0.55),
             fill=LIGHT_BG, line=BORDER)
    add_text(s, "→  VPE 主导阶段一与阶段三，协作阶段四，阶段二保持畅通即可。",
             Inches(0.6), Inches(6.55), Inches(12.13), Inches(0.45),
             size=14, bold=True, color=DARK_MAROON, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, PG[0], TOTAL)


def slide_section_divider(num, name, subtitle):
    s = new_slide()
    add_background(s, DARK_MAROON)
    add_rect(s, Inches(0.6), Inches(2.6), Inches(0.6), Inches(0.6), fill=GOLD)
    add_text(s, f"PHASE {num}  ·  阶段 {num}", Inches(0.6), Inches(3.3), Inches(12), Inches(0.5),
             size=20, bold=True, color=GOLD)
    add_text(s, name, Inches(0.6), Inches(3.85), Inches(12), Inches(1.2),
             size=54, bold=True, color=WHITE)
    add_rect(s, Inches(0.6), Inches(5.2), Inches(2), Emu(28575), fill=GOLD)
    add_text(s, subtitle, Inches(0.6), Inches(5.4), Inches(12), Inches(0.8),
             size=18, color=WARM_LIGHT, italic=True)
    add_text(s, f"{PG[0]} / {TOTAL}", Inches(11.5), Inches(7.2),
             Inches(1.23), Inches(0.25), size=10,
             color=WARM_DIM, align=PP_ALIGN.RIGHT)


def slide_content(step, title, lead, bullets, image=None, image_caption=None):
    s = new_slide()
    add_slide_title(s, f"STEP {step}  ·  步骤 {step}", title, lead)
    add_bullets(s, bullets, Inches(0.6), Inches(2.1), Inches(6.2), Inches(4.6),
                size=15, spacing=10, line_spacing=1.2)
    if image:
        add_image_with_caption(s, image, image_caption or "",
                               Inches(7.0), Inches(2.0),
                               Inches(5.8), Inches(4.2))
    add_footer(s, PG[0], TOTAL)


def slide_text_only(step, title, lead, bullets):
    s = new_slide()
    add_slide_title(s, f"STEP {step}  ·  步骤 {step}", title, lead)
    add_bullets(s, bullets, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5),
                size=16, spacing=10, line_spacing=1.2)
    add_footer(s, PG[0], TOTAL)


# ---------------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------------
slide_cover()
slide_toc()
slide_phase_overview()

# ============ PHASE 1 ============
slide_section_divider(1, "会前规划",
                      "T-7 至 T-1 天  ·  VPE 的主要工作阶段")

slide_content(
    "1.1", "打开议程页面",
    "议程页面是所有会议相关操作的主页。",
    [
        ("路由", "/agenda — 也可通过顶部导航的 Meetings 进入"),
        ("会议选择器", "使用下拉菜单选择要操作的会议"),
        ("状态组", "显示当前会议的状态：Unpublished / Not Started / Running / Finished / Cancelled"),
        ("议程表格", "列出所有环节：时间、角色、负责人、项目、校验标志"),
    ],
    image=os.path.join(SNAP, "agenda_main.png"),
    image_caption="/agenda 页面 — 会议 #982（2026-08-04），状态：Not Started"
)

slide_content(
    "1.2", "从模板创建会议",
    "填写日期、主题与模板，建立新会议。",
    [
        ("点击新建", "在议程编辑器顶部，或通过右上角菜单"),
        ("填写表单", "日期、会议编号、标题/主题、副标题、类型、媒体链接、模板"),
        ("模板文件", "决定默认的议程项目（例如标准俱乐部会议）"),
        ("保存", "POST /agenda/create — 需要 MEETING_CREATE 权限"),
    ],
    image=os.path.join(SNAP, "settings_main.png"),
    image_caption="设置 → 模块 — 生成幻灯片前请确认 'Data/Slides Export' 已开启"
)

slide_content(
    "1.3", "搭建议程",
    "对每个环节进行内联编辑：时间、角色、项目以及会议级字段。",
    [
        ("增删环节", "使用表格底部的行控件"),
        ("调整顺序", "拖动行把手，或使用上下箭头"),
        ("编辑开始时间", "点击时间单元格，输入 24h 制 HH:MM — 冲突会标红"),
        ("每日一词 (WOD)", "在顶部 wod-display 区块设置"),
        ("GE 模式", "ge_mode 下拉框控制总评人报告风格"),
        ("红色徽标", "表示校验问题 — 发布前先解决"),
    ],
    image=os.path.join(SNAP, "agenda_main.png"),
    image_caption="议程表格 — 每一行都可以内联编辑"
)

slide_content(
    "1.4", "为备稿演讲添加项目",
    "备稿演讲环节标记为 Valid_for_Project = true，需要绑定 Pathways 项目。",
    [
        ("点击项目单元格", "在该备稿演讲所在行"),
        ("选择路径", "例如 Dynamic Leadership、Presentation Mastery"),
        ("选择项目", "例如 Level 1 — Ice Breaker"),
        ("可选标题", "会显示在幻灯片和演讲记录中"),
        ("如未定", "留空即可 — 会议启动 (Start) 前补上"),
    ],
    image=os.path.join(SNAP, "pathway_library.png"),
    image_caption="Pathways Library — 备稿演讲项目目录"
)

slide_content(
    "1.5", "开放角色预订",
    "发布 (Publish) 后，所有会员即可预订角色。",
    [
        ("点击 Publish", "在议程页面的状态组中"),
        ("确认提示", "Memory Maker 会提醒你预订将对全体会员开放"),
        ("状态切换", "Unpublished → Not Started"),
        ("会员可见 /booking", "自助预订与候补名单启用"),
        ("提示", "发布后调整议程结构会更麻烦，请先定稿"),
    ],
    image=os.path.join(SNAP, "agenda_main.png"),
    image_caption="议程页面状态组 — 点击 Publish 开放预订"
)

slide_content(
    "1.6", "预订或分配角色",
    "会员自助预订；你负责填补空缺槽位、分配评估人。",
    [
        ("会员自助预订", "访问 /booking，选会议，点击槽位"),
        ("候补名单", "已满的槽位显示「加入候补」 — 负责人取消时自动提升"),
        ("干部分配", "在议程页面点击负责人单元格，选择联系人"),
        ("提前 24–48 小时补位", "填补必需的空白角色；为演讲分配评估人"),
        ("提示", "会前 1–2 天查看 Planner，了解会员的预订意向"),
    ],
    image=os.path.join(SNAP, "booking.png"),
    image_caption="/booking — 会员看到的预订助手页面"
)

slide_text_only(
    "1.7", "会前一天自查清单",
    "在会议前一天逐项过一遍。",
    [
        "议程中没有红色校验徽标",
        "每个备稿演讲都已绑定 Pathways 项目",
        "每个必需角色都有负责人（或有意留空由 TME 现场找人）",
        "每日一词已设置，演讲者已了解词义",
        "幻灯片已生成并保存到笔记本",
        "已与 SAA 确认当晚由谁负责发布和唱票",
        "全部打勾后，你就可以进入阶段三",
    ]
)

# ============ PHASE 2 ============
slide_section_divider(2, "接待签到",
                      "T-0  18:30  ·  由 SAA 主导 — 你只需保持畅通")

slide_content(
    "2.1", "接待与签到",
    "SAA 在门口接待；你只需确认会议已就绪、可以发布。",
    [
        ("路由", "/roster — SAA 在宾客和会员到达时录入"),
        ("票种", "Early-bird / Walk-in / Officer — 顺序号自动计算"),
        ("实时 KPI", "顶部显示总出席人数和收入"),
        ("你的任务", "保持畅通。在 SAA 准备好之前不要发布/启动会议"),
        ("临时编辑", "会议进行中可临时添加即兴演讲者或花名册条目"),
    ],
    image=os.path.join(SNAP, "roster.png"),
    image_caption="/roster — SAA 的签到工具"
)

# ============ PHASE 3 ============
slide_section_divider(3, "现场会议",
                      "T-0  19:00–21:00  ·  状态按钮是你唯一的控制入口")

slide_content(
    "3.1", "发布并启动会议",
    "两次点击：Publish（unpublished → not started），然后 Start（not started → running）。",
    [
        ("Publish", "点击状态组；状态切换为 Not Started"),
        ("等待 18:55", "在正式开场前 5 分钟启动"),
        ("点击 Start", "Start_Time 设为当前时间；状态切换为 Running"),
        ("投票启用", "会员可在 /voting 提交选票；SAA 实时观察票数"),
        ("即时编辑", "会议中可编辑负责人、添加环节、添加即兴演讲者"),
    ],
    image=os.path.join(SNAP, "agenda_main.png"),
    image_caption="状态组：Publish → Start → Stop"
)

slide_text_only(
    "3.2", "应对现场临时变更",
    "常见请求与一分钟内的应对方式。",
    [
        ("添加即兴演讲者", "点击该即兴环节的负责人单元格，选择该会员"),
        ("替换演讲者", "点击负责人单元格，清除原负责人，选择新负责人 — 无需刷新页面"),
        ("临时新增即兴环节", "使用 Add Session 控件；如需调整顺序可拖动"),
        ("忘记分配评估人", "与负责人编辑流程相同 — 从联系人列表中选择即可"),
        ("所有编辑", "通过 POST /agenda/update 一次往返完成，原子化校验和持久化"),
    ]
)

slide_content(
    "3.3", "结束会议",
    "点击 Stop。Memory Maker 封存会议并自动执行收尾工作。",
    [
        ("点击 Stop", "状态切换为 Finished"),
        ("唱票", "获奖者写入投票记录"),
        ("清理候补名单", "该会议的所有候补条目被删除"),
        ("Planner 状态同步", "booked → completed,  waitlist → obsolete,  draft 保持不变"),
        ("项目自动完成", "所有 session log 标记为 Completed，元数据同步"),
    ],
    image=os.path.join(SNAP, "voting.png"),
    image_caption="/voting — 会议停止时自动唱票"
)

# ============ PHASE 4 ============
slide_section_divider(4, "收尾与记录",
                      "T+0 至 T+1  ·  确认唱票、记录演讲、归档会议")

slide_content(
    "4.1", "确认投票结果",
    "SAA 公布结果并运行抽奖；你负责核对数据。",
    [
        ("路由", "/voting — 最终票数与获奖者"),
        ("NPS", "/voting_nps — 俱乐部层面的满意度趋势"),
        ("SAA 核对名单", "确认抽奖对象正确"),
        ("如发现异常", "请勿直接编辑选票 — 通过 /issues 提交问题单以保留审计轨迹"),
    ],
    image=os.path.join(SNAP, "voting.png"),
    image_caption="/voting — Stop 后的最终票数"
)

slide_content(
    "4.2", "更新媒体链接",
    "关联会议录屏，方便会员日后回看。",
    [
        ("回到 /agenda", "在会议头部编辑 Media URL 字段"),
        ("保存", "录屏链接已与会议记录关联"),
        ("出现位置", "会员可在 Planner 和会议档案中查看"),
    ],
    image=os.path.join(SNAP, "contacts.png"),
    image_caption="/contacts — 复核每位演讲者的记录"
)

slide_content(
    "4.3", "复核演讲记录与 Pathways",
    "自动完成通常已处理；只需核对几个关键点。",
    [
        ("/contacts", "确认每位演讲者的项目和级别记录正确"),
        ("自动完成", "会议停止时项目状态已切换为 Completed"),
        ("如项目有误", "通过 /speech_logs 修正并重新保存 — 元数据会重新同步"),
        ("Pathways 进度", "访问 /pathway_library 确认已选路径；级别进度自动递增"),
        ("常见 Bug", "通用占位项目 (ProjectID.GENERIC) 会阻止级别递增 — 替换为真实项目"),
    ],
    image=os.path.join(SNAP, "pathway_library.png"),
    image_caption="Pathways Library — 级别跟踪的项目目录"
)

slide_text_only(
    "4.4", "成就、归档与下一场规划",
    "可选的收尾步骤，以及向下一周期的交接。",
    [
        ("成就", "如俱乐部使用 /achievements，授予新徽章（例如「首次 Ice Breaker」）"),
        ("归档或删除", "默认保留已结束会议在档案中；或点击 Delete 硬删除"),
        ("删除所需权限", "MEETING_CREATE；不可撤销 — 选票、花名册、候补、Planner 条目一并删除"),
        ("规划下一场", "如本次会议设立了新先例，更新模板和设置（参见 template-manager-plan.md）"),
        ("向下一任 VPE 交接", "在任期交接前说明自定义环节、常见问题等"),
    ]
)

# ============ REFERENCE ============
def slide_status_flow():
    s = new_slide()
    add_slide_title(s, "REFERENCE  ·  参考", "会议状态流转",
                    "单向流转。如需重做一场会议，请新建。")

    nodes = [
        ("unpublished", MAROON,     "Unpublished"),
        ("not started", GOLD,       "Not Started"),
        ("running",     DARK_MAROON,"Running"),
        ("finished",    BLUE,       "Finished"),
    ]
    edges = [
        ("Publish", 0, 1),
        ("Start",   1, 2),
        ("Stop",    2, 3),
    ]

    y = Inches(3.2)
    node_w = Inches(2.4)
    node_h = Inches(1.1)
    gap = (SLIDE_W - node_w * 4) / 5
    x_positions = [gap + i * (node_w + gap) for i in range(4)]

    for i, (key, color, label) in enumerate(nodes):
        x = x_positions[i]
        add_rect(s, x, y, node_w, node_h, fill=color, line=color)
        add_text(s, label, x, y, node_w, node_h,
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    for label, src, dst in edges:
        x1 = x_positions[src] + node_w
        x2 = x_positions[dst]
        cy = y + node_h / 2
        add_horizontal_arrow(s, x1, cy, x2)
        add_text(s, label, x1, y - Inches(0.55), x2 - x1, Inches(0.4),
                 size=14, bold=True, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(0.6), Inches(5.5), Inches(12.13), Inches(0.9),
             fill=PANEL_BG, line=BORDER)
    add_text(s,
             "从 Finished → Delete（需 MEETING_CREATE 权限）。硬删除：选票、花名册、"
             "候补名单、Planner 条目和 session log 全部删除。不可撤销。",
             Inches(0.8), Inches(5.55), Inches(11.93), Inches(0.8),
             size=13, color=DARK_TEXT, italic=True, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, "停止时的自动操作：",
             Inches(0.6), Inches(6.5), Inches(12.13), Inches(0.3),
             size=13, bold=True, color=MAROON)
    add_text(s,
             "唱票完成  ·  候补名单清空  ·  Planner 中 booked→completed  ·  "
             "项目自动完成  ·  联系人元数据重新同步",
             Inches(0.6), Inches(6.78), Inches(12.13), Inches(0.3),
             size=12, color=GRAY_TEXT, italic=True)
    add_footer(s, PG[0], TOTAL)


def slide_permissions():
    s = new_slide()
    add_slide_title(s, "REFERENCE  ·  参考", "权限速查表",
                    "本教程默认标准 VPE 角色拥有以下权限。")

    headers = ["操作", "所需权限"]
    rows = [
        ("创建会议",                          "MEETING_CREATE"),
        ("编辑议程（环节、每日一词、项目）",  "MEETING_MANAGE"),
        ("分配 / 移除角色负责人",             "MEETING_MANAGE"),
        ("发布 / 启动 / 结束会议",            "MEETING_MANAGE"),
        ("删除已结束的会议",                  "MEETING_CREATE"),
        ("查看所有会议（含未发布）",          "MEETING_VIEW_ALL"),
        ("生成幻灯片 / 导出",                 "MEETING_MANAGE + 模块 'Data/Slides Export'"),
    ]

    cols = [Inches(7.0), Inches(5.13)]
    row_h = Inches(0.45)
    x0 = Inches(0.6)
    y0 = Inches(2.1)
    add_rect(s, x0, y0, cols[0] + cols[1], row_h, fill=DARK_MAROON)
    cx = x0
    for i, h in enumerate(headers):
        add_text(s, h, cx + Inches(0.15), y0, cols[i], row_h,
                 size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        cx += cols[i]
    for r, (a, p) in enumerate(rows):
        ry = y0 + row_h * (r + 1)
        fill = PANEL_BG if r % 2 == 0 else WHITE
        add_rect(s, x0, ry, cols[0] + cols[1], row_h, fill=fill, line=BORDER, line_w=Pt(0.5))
        cx = x0
        for i, val in enumerate([a, p]):
            color = DARK_TEXT if i == 0 else MAROON
            bold = i == 1
            add_text(s, val, cx + Inches(0.15), ry, cols[i], row_h,
                     size=13, color=color, bold=bold, anchor=MSO_ANCHOR.MIDDLE)
            cx += cols[i]
    add_text(s,
             "如果按钮是灰色，说明缺少对应权限。请联系 ClubAdmin 或在 设置 → 用户 中检查。",
             Inches(0.6), Inches(6.4), Inches(12.13), Inches(0.5),
             size=12, color=GRAY_TEXT, italic=True)
    add_footer(s, PG[0], TOTAL)


def slide_pitfalls():
    s = new_slide()
    add_slide_title(s, "REFERENCE  ·  参考", "常见坑",
                    "前人踩过的雷。提前规避。")

    items = [
        ("过早发布",
         "一旦发布，会员即可预订角色。请先完成结构定稿。"),
        ("绑定错误的 Pathways 项目",
         "错误的项目会污染演讲者记录。如不确定，请询问，不要猜测。"),
        ("误点 Stop",
         "会议会被封存。唯一恢复路径是克隆会议或从备份恢复。"),
        ("忘记分配评估人",
         "会前补上。评估人就是评估类环节的负责人。"),
        ("发布后修改日期",
         "会议编号自动重命名；poster 文件也会重命名。修改后请核对两者。"),
        ("备稿演讲使用了通用项目",
         "ProjectID.GENERIC 会阻止级别递增。替换为真实项目并重新保存。"),
    ]
    add_bullets(s, items, Inches(0.7), Inches(2.1), Inches(11.93), Inches(4.6),
                size=15, spacing=10, line_spacing=1.2)
    add_footer(s, PG[0], TOTAL)


def slide_glossary():
    s = new_slide()
    add_slide_title(s, "REFERENCE  ·  参考", "术语表",
                    "本教程涉及的术语。")

    items = [
        ("VPE",   "教育副主席 — 负责会议内容与 Pathways 进度跟踪。"),
        ("SAA",   "内务副主席 — 负责场地、签到、唱票与抽奖。"),
        ("TME",   "当晚主持人 (Toastmaster of the Evening) — 主理现场会议。"),
        ("WOD",   "每日一词 (Word of the Day) — 传统上由 Memory Maker 负责；演讲者与 TME 使用。"),
        ("Pathway","头马学习路径，由跨多个级别的项目组成。"),
        ("Session Log", "议程表中的一行 — 每个议程项对应一行（TME、即兴、演讲、评估等）。"),
        ("Planner", "会员的个人计划：在哪场会议中完成哪个项目。"),
        ("Waitlist", "系统自动管理的队列；已满槽位的候选人在原负责人取消时自动提升。"),
    ]
    add_bullets(s, items, Inches(0.7), Inches(2.1), Inches(11.93), Inches(4.6),
                size=15, spacing=8, line_spacing=1.15)
    add_footer(s, PG[0], TOTAL)


def slide_where_next():
    s = new_slide()
    add_slide_title(s, "REFERENCE  ·  参考", "延伸阅读",
                    "本仓库中的相关文档。")

    rows = [
        ("SAA — 清单视图",          "docs/SAA_MemMaker_User_Manual.md"),
        ("SAA — 网格视图",          "docs/SAA_User_Manual_Grid.md"),
        ("所有干部速查",            "docs/getting_started.md"),
        ("Pathways 深入",           "docs/ (搜索 'pathways')"),
        ("模板编辑",                "docs/template-manager-plan.md"),
        ("投票 / 实时更新",         "docs/MAILBOX_REALTIME_DESIGN.md"),
        ("用户 / 联系人模型",       "docs/CONTACT_USER_CLUB_MODEL.md"),
        ("本教程的源文档（中文）",  "docs/VPE_MemMaker_Tutorial.zh-CN.md"),
    ]
    y = Inches(2.1)
    row_h = Inches(0.42)
    for i, (label, path) in enumerate(rows):
        ry = y + row_h * i
        fill = PANEL_BG if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.6), ry, Inches(12.13), row_h, fill=fill, line=BORDER, line_w=Pt(0.4))
        add_text(s, label, Inches(0.8), ry, Inches(5.5), row_h,
                 size=14, bold=True, color=DARK_TEXT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, path, Inches(6.4), ry, Inches(6.3), row_h,
                 size=13, color=MAROON, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, PG[0], TOTAL)


def slide_closing():
    s = new_slide()
    add_background(s, DARK_MAROON)
    add_rect(s, Inches(0), Inches(0), Inches(0.5), SLIDE_H, fill=GOLD)
    add_rect(s, Inches(1.5), Inches(2.7), Inches(0.6), Inches(0.6), fill=GOLD)
    add_text(s, "准备就绪", Inches(1.5), Inches(3.4), Inches(11), Inches(0.6),
             size=24, bold=True, color=GOLD)
    add_text(s, "自信地管理会议周期。",
             Inches(1.5), Inches(4.0), Inches(11), Inches(1.0),
             size=40, bold=True, color=WHITE)
    add_text(s,
             "周密规划、发布一次、保持冷静、让 Memory Maker 在结束时自动完成台账。",
             Inches(1.5), Inches(5.0), Inches(11), Inches(1.0),
             size=18, color=WARM_LIGHT, italic=True)
    add_rect(s, Inches(1.5), Inches(6.0), Inches(2), Emu(28575), fill=GOLD)
    add_text(s, "有问题？联系 ClubAdmin 或通过 /issues 提交。",
             Inches(1.5), Inches(6.2), Inches(11), Inches(0.5),
             size=14, color=WARM_DIM, italic=True)
    add_text(s, f"{PG[0]} / {TOTAL}", Inches(11.5), Inches(7.2),
             Inches(1.23), Inches(0.25), size=10,
             color=WARM_DIM, align=PP_ALIGN.RIGHT)


# Build the rest in order
slide_status_flow()
slide_permissions()
slide_pitfalls()
slide_glossary()
slide_where_next()
slide_closing()

# ---------------------------------------------------------------------------
prs.save(OUT)
print(f"Wrote {OUT}  ({PG[0]} slides, {os.path.getsize(OUT)//1024} KB)")
